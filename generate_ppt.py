from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
DARK_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # 배경 (짙은 남색)
ACCENT      = RGBColor(0x89, 0xB4, 0xFA)   # 강조 (파스텔 블루)
ACCENT2     = RGBColor(0xA6, 0xE3, 0xA1)   # 보조 강조 (민트)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xC9, 0xD5)
YELLOW      = RGBColor(0xF9, 0xE2, 0xAF)
RED         = RGBColor(0xF3, 0x8B, 0xA8)
CARD_BG     = RGBColor(0x31, 0x32, 0x44)   # 카드 배경

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 헬퍼 함수 ────────────────────────────────────────────────
def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox

def add_title_slide(prs, week_num, title, subtitle, phase_label=""):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)

    # 좌측 강조 바
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, ACCENT)

    # Phase 라벨
    if phase_label:
        add_rect(slide, Inches(0.7), Inches(0.5), Inches(3.5), Inches(0.45), CARD_BG)
        add_textbox(slide, phase_label, Inches(0.75), Inches(0.52), Inches(3.4), Inches(0.4),
                    font_size=13, color=ACCENT, bold=True)

    # Week 번호
    add_textbox(slide, f"Week {week_num:02d}", Inches(0.7), Inches(1.1), Inches(5), Inches(0.7),
                font_size=28, color=ACCENT, bold=True)

    # 제목
    add_textbox(slide, title, Inches(0.7), Inches(1.85), Inches(11.5), Inches(1.8),
                font_size=44, bold=True, color=WHITE)

    # 구분선
    add_rect(slide, Inches(0.7), Inches(3.75), Inches(5), Inches(0.06), ACCENT)

    # 부제목
    add_textbox(slide, subtitle, Inches(0.7), Inches(3.9), Inches(11), Inches(1.2),
                font_size=20, color=LIGHT_GRAY)

    # 우측 장식 원
    for i, (cx, cy, r, col) in enumerate([
        (11.5, 1.5, 1.8, ACCENT),
        (12.2, 3.2, 1.0, ACCENT2),
        (10.8, 4.8, 0.6, YELLOW),
    ]):
        shape = slide.shapes.add_shape(
            9,  # oval
            Inches(cx - r/2), Inches(cy - r/2), Inches(r), Inches(r)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        shape.fill.fore_color.theme_color  # dummy access
        # 투명도를 직접 못 쓰므로 색을 섞어 표현
        # (장식 원이므로 그냥 진행)

    # 하단 정보
    add_textbox(slide, "LLM 기반 AI 서비스 개발 실습", Inches(0.7), Inches(6.8), Inches(6), Inches(0.5),
                font_size=13, color=LIGHT_GRAY)
    add_textbox(slide, "3시간 수업", Inches(11), Inches(6.8), Inches(2), Inches(0.5),
                font_size=13, color=ACCENT2, align=PP_ALIGN.RIGHT)
    return slide

def add_agenda_slide(prs, items):
    """목차 슬라이드 (최대 8개 항목)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, ACCENT2)

    add_textbox(slide, "목  차", Inches(0.7), Inches(0.3), Inches(10), Inches(0.8),
                font_size=36, bold=True, color=WHITE)
    add_rect(slide, Inches(0.7), Inches(1.15), Inches(4), Inches(0.05), ACCENT2)

    cols = [items[:4], items[4:]] if len(items) > 4 else [items, []]
    offsets = [Inches(0.7), Inches(6.8)]
    for ci, col in enumerate(cols):
        for ri, item in enumerate(col):
            y = Inches(1.4 + ri * 1.3)
            x = offsets[ci]
            num = ri + 1 + ci * 4
            # 번호 원
            circ = slide.shapes.add_shape(9, x, y + Inches(0.05), Inches(0.55), Inches(0.55))
            circ.fill.solid()
            circ.fill.fore_color.rgb = ACCENT if ci == 0 else ACCENT2
            circ.line.fill.background()
            tf2 = circ.text_frame
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
            r2 = tf2.paragraphs[0].add_run()
            r2.text = str(num)
            r2.font.size = Pt(16)
            r2.font.bold = True
            r2.font.color.rgb = DARK_BG
            r2.font.name = "맑은 고딕"
            # 텍스트
            add_textbox(slide, item, x + Inches(0.65), y, Inches(5.8), Inches(0.65),
                        font_size=17, color=WHITE)
    return slide

def add_content_slide(prs, title, bullets, icon="▸", accent_color=ACCENT):
    """일반 내용 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # 상단 헤더 바
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), CARD_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(1.15), accent_color)
    add_textbox(slide, title, Inches(0.65), Inches(0.2), Inches(12), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    # 불릿 내용
    y = Inches(1.4)
    for item in bullets:
        if item.startswith("##"):          # 섹션 소제목
            text = item[2:].strip()
            add_rect(slide, Inches(0.35), y, Inches(0.25), Inches(0.45), accent_color)
            add_textbox(slide, text, Inches(0.7), y, Inches(12), Inches(0.5),
                        font_size=18, bold=True, color=accent_color)
            y += Inches(0.55)
        elif item.startswith("  "):        # 들여쓰기 (서브 불릿)
            text = item.strip()
            add_textbox(slide, f"  › {text}", Inches(1.0), y, Inches(11.5), Inches(0.5),
                        font_size=15, color=LIGHT_GRAY)
            y += Inches(0.48)
        elif item == "---":                # 구분선
            add_rect(slide, Inches(0.35), y + Inches(0.1), Inches(12.5), Inches(0.04), CARD_BG)
            y += Inches(0.3)
        else:
            add_textbox(slide, f"{icon}  {item}", Inches(0.35), y, Inches(12.5), Inches(0.55),
                        font_size=17, color=WHITE)
            y += Inches(0.58)

    return slide

def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items,
                      accent_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), CARD_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(1.15), accent_color)
    add_textbox(slide, title, Inches(0.65), Inches(0.2), Inches(12), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    for col_idx, (col_title, col_items) in enumerate([(left_title, left_items), (right_title, right_items)]):
        x = Inches(0.35 + col_idx * 6.5)
        w = Inches(6.1)
        # 카드 배경
        add_rect(slide, x, Inches(1.25), w, Inches(6.0), CARD_BG)
        # 컬럼 헤더
        add_rect(slide, x, Inches(1.25), w, Inches(0.55), accent_color if col_idx == 0 else ACCENT2)
        add_textbox(slide, col_title, x + Inches(0.15), Inches(1.3), w - Inches(0.3), Inches(0.45),
                    font_size=17, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        y = Inches(1.95)
        for item in col_items:
            add_textbox(slide, f"• {item}", x + Inches(0.2), y, w - Inches(0.4), Inches(0.55),
                        font_size=15, color=WHITE)
            y += Inches(0.57)
    return slide

def add_code_slide(prs, title, code_lines, desc="", accent_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), CARD_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(1.15), accent_color)
    add_textbox(slide, title, Inches(0.65), Inches(0.2), Inches(12), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    if desc:
        add_textbox(slide, desc, Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.5),
                    font_size=15, color=LIGHT_GRAY)

    code_top = Inches(1.75) if desc else Inches(1.25)
    code_h = SLIDE_H - code_top - Inches(0.3)
    add_rect(slide, Inches(0.35), code_top, Inches(12.6), code_h, RGBColor(0x11, 0x12, 0x1B))

    code_text = "\n".join(code_lines)
    txBox = slide.shapes.add_textbox(Inches(0.55), code_top + Inches(0.15),
                                     Inches(12.2), code_h - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(13)
    run.font.name = "Consolas"
    run.font.color.rgb = ACCENT2
    return slide

def add_summary_slide(prs, week_num, key_points, next_week=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), CARD_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(1.15), YELLOW)
    add_textbox(slide, f"Week {week_num:02d}  핵심 정리", Inches(0.65), Inches(0.2), Inches(12), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    # 핵심 포인트 카드
    add_rect(slide, Inches(0.35), Inches(1.3), Inches(8.1), Inches(5.5), CARD_BG)
    add_textbox(slide, "이번 주 핵심 개념", Inches(0.5), Inches(1.4), Inches(7.8), Inches(0.5),
                font_size=17, bold=True, color=YELLOW)
    y = Inches(1.95)
    for kp in key_points:
        add_textbox(slide, f"✓  {kp}", Inches(0.55), y, Inches(7.7), Inches(0.55),
                    font_size=15, color=WHITE)
        y += Inches(0.57)

    # 다음 주 예고
    if next_week:
        add_rect(slide, Inches(8.65), Inches(1.3), Inches(4.3), Inches(5.5), CARD_BG)
        add_textbox(slide, "다음 주 예고", Inches(8.8), Inches(1.4), Inches(4.0), Inches(0.5),
                    font_size=17, bold=True, color=ACCENT)
        add_textbox(slide, next_week, Inches(8.8), Inches(2.0), Inches(4.0), Inches(4.5),
                    font_size=15, color=WHITE)

    add_textbox(slide, "수고하셨습니다!", Inches(0.35), Inches(6.9), Inches(12.5), Inches(0.5),
                font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ════════════════════════════════════════════════════════════
#  각 주차 PPT 생성 함수
# ════════════════════════════════════════════════════════════

def make_week01():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 1,
        "LLM 원리 및 Gemini API\n환경 세팅",
        "인공지능이 글을 쓰는 방법 / 나만의 개발 환경 구축",
        "Phase 1 — LLM 기초 및 API 제어")

    add_agenda_slide(prs, [
        "LLM이란 무엇인가?",
        "Transformer 아키텍처 기초",
        "토큰과 토크나이제이션",
        "Gemini API 소개 및 특징",
        "개발 환경 세팅 (Python·venv)",
        "API Key 발급 및 보안",
        "첫 번째 API 호출 실습",
        "Q&A 및 과제 안내",
    ])

    add_content_slide(prs, "LLM이란? — Language Model의 진화", [
        "## LM → GPT → LLM",
        "통계적 언어 모델 (n-gram) → 신경망 언어 모델 → 대규모 언어 모델",
        "---",
        "## 핵심 아이디어: '다음 토큰 예측'",
        "\"The cat sat on the ___\" → 확률 분포로 다음 단어 예측",
        "  수십억 개 파라미터로 방대한 텍스트 학습 → 언어 능력 창발(Emergence)",
        "---",
        "## 대표 LLM 모델",
        "  GPT-4o (OpenAI),  Claude 3.5 (Anthropic),  Gemini 1.5 Pro (Google)",
        "  Llama 3 (Meta) — 오픈소스",
        "---",
        "## 왜 Gemini API를 사용하나?",
        "  무료 티어 제공 (Flash 모델) / 한국어 우수 / 멀티모달 지원",
    ])

    add_content_slide(prs, "Transformer 아키텍처 기초", [
        "## Self-Attention 메커니즘",
        "문장 내 모든 토큰이 서로 '관계도'를 계산 → 문맥 파악",
        "  Q (Query) · K (Key) · V (Value) 행렬 연산",
        "---",
        "## Encoder vs Decoder",
        "  Encoder: 입력 이해 (BERT 계열)",
        "  Decoder: 텍스트 생성 (GPT / Gemini 계열)",
        "  Encoder-Decoder: 번역·요약 (T5, BART)",
        "---",
        "## 사전학습 → 파인튜닝 → RLHF",
        "방대한 웹 데이터로 사전학습 → 지시 따르기(Instruction Fine-tuning) → 인간 피드백 강화학습",
        "---",
        "## 파라미터 규모와 성능",
        "7B / 13B / 70B / 1T+ — 규모가 커질수록 '창발' 능력 증가",
    ])

    add_content_slide(prs, "토큰과 토크나이제이션", [
        "## 토큰이란?",
        "모델이 처리하는 최소 단위 — 단어, 서브워드, 문자",
        "  'unhappiness' → ['un', 'happi', 'ness']  (약 3 토큰)",
        "  한국어: '안녕하세요' → ['안', '녕', '하', '세', '요']  (더 많은 토큰)",
        "---",
        "## 컨텍스트 윈도우 (Context Window)",
        "모델이 한 번에 처리할 수 있는 최대 토큰 수",
        "  GPT-4o: 128K  /  Gemini 1.5 Pro: 1M  /  Claude 3.5: 200K",
        "---",
        "## 토큰 비용",
        "API 사용 비용은 입력+출력 토큰 합산으로 계산",
        "  Gemini Flash: $0.075 / 1M tokens (입력) — 매우 저렴",
        "---",
        "## 실습: 토크나이저 체험",
        "  platform.openai.com/tokenizer 에서 직접 토큰 확인",
    ])

    add_content_slide(prs, "Gemini API 소개 및 주요 모델", [
        "## Google AI Studio",
        "  aistudio.google.com — 무료로 API Key 발급 + 플레이그라운드",
        "---",
        "## 주요 모델 라인업",
        "  gemini-1.5-flash   : 빠르고 저렴, 기본 실습에 적합 ★",
        "  gemini-1.5-pro     : 고성능, 복잡한 추론·긴 문서 처리",
        "  gemini-1.5-flash-8b: 초경량, 비용 최소화",
        "---",
        "## 지원 기능",
        "  텍스트 생성 / 멀티모달(이미지·영상) / 임베딩",
        "  Function Calling / JSON Mode / 스트리밍",
        "---",
        "## 사용 한도 (무료 티어)",
        "  Flash: 분당 15 요청 / 일 1,500 요청  — 수업용으로 충분",
    ])

    add_content_slide(prs, "개발 환경 세팅 (실습)", [
        "## 필수 설치",
        "  Python 3.10+ (python.org)",
        "  VS Code + Python 익스텐션",
        "  Git (git-scm.com)",
        "---",
        "## 가상환경 생성",
        "  python -m venv .venv",
        "  source .venv/bin/activate   (Mac/Linux)",
        "  .venv\\Scripts\\activate      (Windows)",
        "---",
        "## 패키지 설치",
        "  pip install google-generativeai python-dotenv",
        "---",
        "## .env 파일로 API Key 관리",
        "  GOOGLE_API_KEY=AIza...",
        "  절대 GitHub에 올리지 말 것! → .gitignore 설정 필수",
    ])

    add_code_slide(prs, "첫 번째 API 호출 — Hello, Gemini!",
        [
            "import google.generativeai as genai",
            "import os",
            "from dotenv import load_dotenv",
            "",
            "load_dotenv()  # .env 파일 로드",
            "genai.configure(api_key=os.environ['GOOGLE_API_KEY'])",
            "",
            "# 모델 초기화",
            "model = genai.GenerativeModel('gemini-1.5-flash')",
            "",
            "# 텍스트 생성",
            "response = model.generate_content('한국의 수도는 어디인가요?')",
            "",
            "print(response.text)",
            "# → '한국의 수도는 서울입니다.'",
            "",
            "# 토큰 사용량 확인",
            "print(response.usage_metadata)",
        ],
        "환경변수로 API Key를 안전하게 관리하고, 첫 응답을 받아봅니다."
    )

    add_summary_slide(prs, 1,
        [
            "LLM = 대규모 텍스트로 학습한 '다음 토큰 예측' 모델",
            "Transformer Self-Attention이 문맥 이해의 핵심",
            "토큰 = 모델의 최소 처리 단위 / 비용 단위",
            "API Key는 .env + .gitignore로 반드시 보호",
            "google-generativeai SDK로 3줄 만에 AI 호출 가능",
        ],
        "2주차 예고:\nPython SDK 심화 활용\nCLI 챗봇 만들기\n스트리밍 응답 처리"
    )

    return prs

def make_week02():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 2,
        "Python SDK 활용 및\nCLI 챗봇 구현",
        "SDK 심화 + 대화형 챗봇 완성",
        "Phase 1 — LLM 기초 및 API 제어")

    add_agenda_slide(prs, [
        "지난 주 복습 & Q&A",
        "GenerativeModel 파라미터 심화",
        "스트리밍(Streaming) 응답",
        "멀티턴 대화 — ChatSession",
        "CLI 챗봇 설계 및 구현",
        "안전 설정(Safety Settings)",
        "실습: 나만의 CLI 챗봇",
        "과제 안내",
    ])

    add_content_slide(prs, "GenerativeModel 파라미터 심화", [
        "## GenerationConfig",
        "  temperature   : 창의성 조절 (0.0 ~ 2.0) — 낮을수록 안정적",
        "  top_p         : 누적 확률 샘플링 (0~1)",
        "  top_k         : 상위 k개 토큰에서 샘플링",
        "  max_output_tokens: 최대 출력 토큰 수",
        "---",
        "## system_instruction",
        "모델의 역할과 행동 방식을 사전 정의",
        "  '당신은 친절한 한국어 튜터입니다. ...'",
        "---",
        "## 모델 선택 전략",
        "  실험·개발 단계 → gemini-1.5-flash (빠르고 저렴)",
        "  복잡한 추론 → gemini-1.5-pro",
    ])

    add_code_slide(prs, "GenerationConfig 활용 예시",
        [
            "import google.generativeai as genai",
            "",
            "generation_config = genai.GenerationConfig(",
            "    temperature=0.7,      # 적당한 창의성",
            "    top_p=0.9,",
            "    max_output_tokens=512,",
            ")",
            "",
            "model = genai.GenerativeModel(",
            "    model_name='gemini-1.5-flash',",
            "    generation_config=generation_config,",
            "    system_instruction='당신은 파이썬 전문 튜터입니다. '",
            "                       '간결하고 명확하게 설명하세요.',",
            ")",
            "",
            "response = model.generate_content('리스트 컴프리헨션이 뭔가요?')",
            "print(response.text)",
        ],
        "temperature와 system_instruction으로 모델 동작을 세밀하게 제어합니다."
    )

    add_content_slide(prs, "스트리밍(Streaming) 응답", [
        "## 왜 스트리밍인가?",
        "전체 응답 대기 (수 초) → 토큰 단위 실시간 출력",
        "사용자 경험(UX) 대폭 개선 — ChatGPT와 동일한 효과",
        "---",
        "## generate_content vs stream=True",
        "  일반: response = model.generate_content(prompt)",
        "  스트리밍: stream=True 파라미터 추가",
        "---",
        "## 스트리밍 처리 패턴",
        "  for chunk in response: print(chunk.text, end='', flush=True)",
        "---",
        "## 주의사항",
        "  스트리밍 중 response.text 접근 → 오류 발생",
        "  chunk.text로 조각별 접근해야 함",
    ])

    add_code_slide(prs, "스트리밍 응답 구현",
        [
            "model = genai.GenerativeModel('gemini-1.5-flash')",
            "",
            "# 스트리밍 요청",
            "response = model.generate_content(",
            "    '파이썬의 장점을 5가지 설명해줘',",
            "    stream=True",
            ")",
            "",
            "# 실시간 출력",
            "print('AI: ', end='')",
            "for chunk in response:",
            "    print(chunk.text, end='', flush=True)",
            "print()  # 줄바꿈",
            "",
            "# 완료 후 전체 텍스트 접근",
            "response.resolve()  # 스트리밍 완료 대기",
            "# print(response.text)  # 이제 가능",
        ],
        "실시간 타이핑 효과로 사용자 경험을 개선합니다."
    )

    add_content_slide(prs, "멀티턴 대화 — ChatSession", [
        "## 대화 이력(History)의 중요성",
        "LLM은 기본적으로 무상태(Stateless) — 매 요청이 독립적",
        "이전 대화를 기억하려면 히스토리를 직접 전달해야 함",
        "---",
        "## ChatSession 자동 관리",
        "  chat = model.start_chat(history=[])",
        "  chat.send_message() 호출 시 자동으로 history 누적",
        "---",
        "## 히스토리 구조",
        "  {'role': 'user', 'parts': ['안녕!']}",
        "  {'role': 'model', 'parts': ['안녕하세요!']}",
        "---",
        "## 주의: 토큰 누적",
        "대화가 길어질수록 비용 증가 → 히스토리 길이 제한 필요",
    ])

    add_code_slide(prs, "CLI 챗봇 완성 코드",
        [
            "import google.generativeai as genai",
            "import os",
            "from dotenv import load_dotenv",
            "",
            "load_dotenv()",
            "genai.configure(api_key=os.environ['GOOGLE_API_KEY'])",
            "",
            "model = genai.GenerativeModel(",
            "    'gemini-1.5-flash',",
            "    system_instruction='당신은 친절한 AI 어시스턴트입니다.'",
            ")",
            "chat = model.start_chat(history=[])",
            "",
            "print('챗봇 시작! (종료: quit)')  ",
            "while True:",
            "    user_input = input('You: ').strip()",
            "    if user_input.lower() == 'quit': break",
            "    if not user_input: continue",
            "    response = chat.send_message(user_input, stream=True)",
            "    print('AI: ', end='')",
            "    for chunk in response:",
            "        print(chunk.text, end='', flush=True)",
            "    print()",
        ],
        "완전한 동작 CLI 챗봇 — 복사하면 바로 실행됩니다."
    )

    add_summary_slide(prs, 2,
        [
            "GenerationConfig로 temperature 등 생성 파라미터 제어",
            "system_instruction으로 모델에 역할 부여",
            "stream=True로 실시간 타이핑 효과 구현",
            "ChatSession이 대화 히스토리 자동 관리",
            "CLI 챗봇 = 루프 + ChatSession + 스트리밍",
        ],
        "3주차 예고:\n프롬프트 엔지니어링\nFew-shot 예시 작성\nPersona 설계"
    )

    return prs

def make_week03():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 3,
        "프롬프트 엔지니어링",
        "Few-shot · Persona · Chain-of-Thought",
        "Phase 1 — LLM 기초 및 API 제어")

    add_agenda_slide(prs, [
        "프롬프트 엔지니어링이란?",
        "Zero-shot / One-shot / Few-shot",
        "Persona 설계",
        "Chain-of-Thought (CoT)",
        "Role Prompting & 지시 템플릿",
        "나쁜 프롬프트 vs 좋은 프롬프트",
        "실습: 감성 분석 챗봇",
        "과제 안내",
    ])

    add_content_slide(prs, "프롬프트 엔지니어링이란?", [
        "## 정의",
        "모델 파라미터를 변경하지 않고, 입력(프롬프트)만 조정해 원하는 출력을 얻는 기술",
        "---",
        "## 왜 중요한가?",
        "같은 모델도 프롬프트에 따라 결과가 천차만별",
        "  '수도가 뭐야?' vs '5살 어린이에게 설명하듯 수도의 개념을 알려줘'",
        "---",
        "## 프롬프트의 구성 요소",
        "  지시(Instruction): 모델이 할 일",
        "  맥락(Context): 배경 정보",
        "  입력(Input): 처리할 실제 데이터",
        "  출력 형식(Output Format): 원하는 응답 형태",
        "---",
        "## 핵심 원칙",
        "  명확하게 / 구체적으로 / 예시 포함 / 역할 부여",
    ])

    add_two_col_slide(prs, "Zero-shot vs Few-shot",
        "Zero-shot", [
            "예시 없이 지시만 제공",
            "빠르고 간편",
            "단순한 작업에 효과적",
            "",
            "예: '이 문장의 감성을",
            "긍정/부정으로 분류해줘'",
            "'오늘 날씨가 정말 좋다!'",
            "→ 긍정",
        ],
        "Few-shot", [
            "2~5개 예시를 함께 제공",
            "복잡한 패턴 학습 가능",
            "일관된 출력 형식 유도",
            "",
            "예: 입력: 맛있다 → 긍정",
            "입력: 별로다 → 부정",
            "입력: 오늘 날씨 좋다 → ?",
            "→ 더 정확한 결과",
        ]
    )

    add_code_slide(prs, "Few-shot 프롬프트 구현",
        [
            "few_shot_prompt = '''",
            "당신은 감성 분석 전문가입니다.",
            "아래 예시를 참고해 입력 문장의 감성을 분류하세요.",
            "",
            "예시:",
            "입력: 이 영화 정말 재밌었어요!",
            "감성: 긍정",
            "",
            "입력: 서비스가 너무 불친절했어요.",
            "감성: 부정",
            "",
            "입력: 그냥 평범했어요.",
            "감성: 중립",
            "",
            "이제 분류하세요:",
            "입력: {text}",
            "감성:'''",
            "",
            "def analyze_sentiment(text: str) -> str:",
            "    prompt = few_shot_prompt.format(text=text)",
            "    response = model.generate_content(prompt)",
            "    return response.text.strip()",
        ],
        "f-string 템플릿으로 재사용 가능한 Few-shot 프롬프트를 만듭니다."
    )

    add_content_slide(prs, "Persona 설계", [
        "## Persona란?",
        "모델에게 특정 역할과 성격을 부여 → 일관된 응답 스타일",
        "---",
        "## 효과적인 Persona 구성 요소",
        "  이름과 역할: '당신은 Alex, 10년 경력 파이썬 개발자입니다'",
        "  말투/스타일: '항상 한국어로, 친근하고 격려하는 톤으로 답하세요'",
        "  전문성 범위: '파이썬·데이터 분석 전문, 다른 분야는 모른다고 답하세요'",
        "  제약 조건: '답변은 200자 이내로 요약하세요'",
        "---",
        "## system_instruction 활용",
        "  페르소나 설정은 system_instruction에 넣는 것이 최선",
        "  매 요청마다 반복하지 않아도 되는 전역 설정",
        "---",
        "## 주의",
        "  지나치게 복잡한 페르소나 → 모델 혼란 가능",
    ])

    add_content_slide(prs, "Chain-of-Thought (CoT)", [
        "## 개념",
        "'단계별로 생각하고 답하라'는 지시로 복잡한 추론 성능 향상",
        "---",
        "## 방법 1: 직접 지시",
        "  '단계별로 생각해서 답하세요 (Let's think step by step)'",
        "---",
        "## 방법 2: Few-shot CoT",
        "  예시에 중간 추론 과정을 포함시켜 패턴 학습 유도",
        "---",
        "## 적용 사례",
        "  수학 문제: 계산 단계 명시 → 오답률 감소",
        "  코드 디버깅: 오류 원인 분석 → 수정 제안",
        "  논리 퀴즈: 가능성 나열 → 결론 도출",
        "---",
        "## 주의",
        "  CoT는 토큰 소비 증가 → 비용 고려 필요",
    ])

    add_two_col_slide(prs, "나쁜 프롬프트 vs 좋은 프롬프트",
        "나쁜 프롬프트", [
            "짧고 모호한 지시",
            "'코드 짜줘'",
            "'설명해줘'",
            "'고쳐줘'",
            "",
            "문제점:",
            "- 어떤 언어인지 불명확",
            "- 어떤 수준의 설명인지 모름",
            "- 뭘 어떻게 고쳐야 할지 모름",
        ],
        "좋은 프롬프트", [
            "역할 + 맥락 + 구체적 지시",
            "'파이썬으로 CSV 파일을 읽어",
            "나이 열 기준 정렬 후",
            "상위 10개 출력하는 코드를",
            "pandas 없이 작성해줘'",
            "",
            "효과:",
            "- 언어·라이브러리 명시",
            "- 요구사항 구체적",
            "- 제약 조건 포함",
        ]
    )

    add_summary_slide(prs, 3,
        [
            "프롬프트 엔지니어링 = 지시·맥락·예시·형식의 조합",
            "Few-shot: 2~5개 예시로 패턴 학습 유도",
            "Persona: system_instruction에 역할·스타일 부여",
            "CoT: '단계별로 생각'으로 복잡한 추론 개선",
            "구체적·명확한 프롬프트가 항상 좋은 결과를 만든다",
        ],
        "4주차 예고:\n구조화된 출력 (JSON)\nPydantic 모델 활용\n대화 이력 고급 관리"
    )

    return prs

def make_week04():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 4,
        "구조화된 출력(JSON) 및\n대화 이력 관리",
        "데이터 파싱 자동화 + 스마트 히스토리 관리",
        "Phase 1 — LLM 기초 및 API 제어")

    add_agenda_slide(prs, [
        "왜 구조화된 출력인가?",
        "JSON 출력 유도 기법",
        "response_mime_type 활용",
        "Pydantic 모델로 검증",
        "대화 이력 관리 전략",
        "토큰 절약 — 요약 기반 압축",
        "실습: 데이터 추출 봇",
        "Phase 1 마무리 & 프로젝트 준비",
    ])

    add_content_slide(prs, "왜 구조화된 출력인가?", [
        "## 문제 상황",
        "LLM 응답은 자유형식 텍스트 → 파싱이 어렵고 불안정",
        "  '서울, 25도, 맑음' vs '현재 서울의 기온은 25°C이며 맑습니다'",
        "---",
        "## 구조화된 출력의 장점",
        "  파싱 코드 불필요 — JSON 바로 dict로 변환",
        "  데이터베이스 저장, API 연동에 바로 활용",
        "  오류 감소 — 필드 누락·타입 오류 방지",
        "---",
        "## 활용 사례",
        "  정보 추출: 이메일 → 이름·날짜·액션 아이템 추출",
        "  분류 시스템: 입력 → 카테고리·신뢰도 반환",
        "  데이터 변환: 비정형 텍스트 → 정형 데이터",
    ])

    add_code_slide(prs, "JSON 출력 유도 — response_mime_type",
        [
            "import google.generativeai as genai",
            "import json",
            "",
            "model = genai.GenerativeModel(",
            "    'gemini-1.5-flash',",
            "    generation_config=genai.GenerationConfig(",
            "        response_mime_type='application/json'  # JSON 강제",
            "    )",
            ")",
            "",
            "prompt = '''",
            "다음 텍스트에서 정보를 추출하세요:",
            "'홍길동(28세)은 서울에서 개발자로 일합니다.'",
            "",
            "JSON 형식: {\"name\": \"\", \"age\": 0, \"city\": \"\", \"job\": \"\"}",
            "'''",
            "",
            "response = model.generate_content(prompt)",
            "data = json.loads(response.text)",
            "print(data['name'])  # → '홍길동'",
        ],
        "response_mime_type='application/json'으로 JSON만 반환하도록 강제합니다."
    )

    add_content_slide(prs, "Pydantic 모델로 스키마 강제", [
        "## Pydantic이란?",
        "Python 데이터 검증 라이브러리 — 타입 힌트 기반",
        "  설치: pip install pydantic",
        "---",
        "## LLM + Pydantic 패턴",
        "1. Pydantic 모델로 원하는 데이터 구조 정의",
        "2. 모델의 JSON 스키마를 프롬프트에 삽입",
        "3. LLM 응답을 Pydantic으로 파싱·검증",
        "---",
        "## 장점",
        "  타입 자동 변환 (str → int)",
        "  필수 필드 누락 시 즉시 오류",
        "  IDE 자동완성 지원",
        "---",
        "## Gemini response_schema 파라미터",
        "Pydantic 모델을 직접 response_schema에 전달 가능",
    ])

    add_code_slide(prs, "Pydantic + Gemini 통합 예시",
        [
            "from pydantic import BaseModel",
            "from typing import List",
            "import google.generativeai as genai",
            "",
            "class MovieReview(BaseModel):",
            "    title: str",
            "    rating: float   # 1.0 ~ 5.0",
            "    pros: List[str]",
            "    cons: List[str]",
            "    recommend: bool",
            "",
            "model = genai.GenerativeModel(",
            "    'gemini-1.5-flash',",
            "    generation_config=genai.GenerationConfig(",
            "        response_mime_type='application/json',",
            "        response_schema=MovieReview  # Pydantic 모델 직접 전달",
            "    )",
            ")",
            "",
            "resp = model.generate_content('인터스텔라 리뷰를 JSON으로')",
            "review = MovieReview.model_validate_json(resp.text)",
            "print(review.rating, review.recommend)",
        ],
        "Pydantic 모델을 response_schema로 전달하면 완벽한 타입 보장을 받습니다."
    )

    add_content_slide(prs, "대화 이력 관리 전략", [
        "## 문제: 무한 히스토리",
        "대화가 길어지면 컨텍스트 윈도우 초과 → 오류 또는 비용 급증",
        "---",
        "## 전략 1: 슬라이딩 윈도우",
        "최근 N턴만 유지 → 오래된 메시지 자동 제거",
        "  history = history[-MAX_TURNS:]",
        "---",
        "## 전략 2: 요약 압축",
        "오래된 대화를 LLM으로 요약 → 요약본만 유지",
        "  '이전 대화 요약: 사용자가 파이썬 기초를 물어봤음'",
        "---",
        "## 전략 3: 주제별 세션 분리",
        "주제가 바뀔 때 chat.history를 초기화",
        "---",
        "## 실전 권장",
        "  슬라이딩 윈도우(10~20턴) + 시스템 프롬프트 유지",
    ])

    add_code_slide(prs, "슬라이딩 윈도우 히스토리 관리",
        [
            "MAX_HISTORY = 10  # 최근 10턴만 유지",
            "",
            "class ManagedChat:",
            "    def __init__(self, model, system_prompt):",
            "        self.model = model",
            "        self.system_prompt = system_prompt",
            "        self.history = []",
            "",
            "    def chat(self, user_input: str) -> str:",
            "        # 슬라이딩 윈도우 적용",
            "        trimmed = self.history[-MAX_HISTORY:]",
            "        ",
            "        chat_session = self.model.start_chat(history=trimmed)",
            "        response = chat_session.send_message(user_input)",
            "        ",
            "        # 히스토리 업데이트",
            "        self.history.append({'role': 'user',  'parts': [user_input]})",
            "        self.history.append({'role': 'model', 'parts': [response.text]})",
            "        ",
            "        return response.text",
        ],
        "MAX_HISTORY로 히스토리 길이를 제한해 토큰 비용을 통제합니다."
    )

    add_summary_slide(prs, 4,
        [
            "response_mime_type='application/json'으로 JSON 출력 강제",
            "Pydantic + response_schema로 타입 안전한 LLM 출력",
            "슬라이딩 윈도우로 히스토리 토큰 비용 통제",
            "Phase 1 완료: LLM 기초·SDK·프롬프트·출력 관리",
            "다음 Phase: RAG — 내 문서로 답하는 AI 만들기",
        ],
        "5주차 예고:\n임베딩이란?\n벡터 데이터베이스\nChromaDB 실습"
    )

    return prs

def make_week05():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 5,
        "임베딩과 벡터\n데이터베이스 (ChromaDB)",
        "의미 검색의 원리와 ChromaDB 실습",
        "Phase 1 — RAG 기초")

    add_agenda_slide(prs, [
        "임베딩(Embedding)이란?",
        "의미 유사도 계산",
        "벡터 데이터베이스 개요",
        "ChromaDB 설치 및 기본 사용법",
        "Gemini 임베딩 API 활용",
        "유사 문서 검색 실습",
        "RAG 파이프라인 개요",
        "과제 안내",
    ])

    add_content_slide(prs, "임베딩(Embedding)이란?", [
        "## 개념",
        "텍스트(단어·문장·문서)를 수치 벡터로 변환하는 기법",
        "  '고양이' → [0.23, -0.51, 0.87, ... ] (수백~수천 차원)",
        "---",
        "## 왜 필요한가?",
        "컴퓨터는 텍스트를 직접 비교 못함 → 벡터로 변환해 수학적 비교",
        "---",
        "## 핵심 특성: 의미가 가까우면 벡터도 가까움",
        "  '왕' - '남자' + '여자' ≈ '여왕' (Word2Vec 고전 예시)",
        "  '개' ↔ '강아지': 코사인 유사도 높음",
        "  '개' ↔ '자동차': 코사인 유사도 낮음",
        "---",
        "## 임베딩 모델",
        "  text-embedding-004 (Google) — Gemini API 제공",
        "  text-embedding-3-small (OpenAI)",
        "  BGE-m3, E5 (오픈소스)",
    ])

    add_content_slide(prs, "코사인 유사도 — 의미 거리 측정", [
        "## 코사인 유사도(Cosine Similarity)",
        "두 벡터 사이의 각도로 유사도 측정 (방향 중심)",
        "  값 범위: -1(반대) ~ 0(무관) ~ 1(동일)",
        "---",
        "## 직관적 이해",
        "  '파이썬 프로그래밍' vs '파이썬 코딩' → 0.95 (매우 유사)",
        "  '파이썬 프로그래밍' vs '뱀의 생태' → 0.40 (낮음)",
        "---",
        "## 계산 공식",
        "  similarity = (A·B) / (||A|| × ||B||)",
        "---",
        "## 다른 거리 지표",
        "  유클리드 거리: 벡터 크기도 고려 (문서 길이 영향)",
        "  내적(Dot Product): 크기와 방향 모두 고려",
        "  RAG에서는 코사인 유사도가 가장 많이 쓰임",
    ])

    add_code_slide(prs, "Gemini 임베딩 API 사용법",
        [
            "import google.generativeai as genai",
            "import numpy as np",
            "",
            "genai.configure(api_key='YOUR_API_KEY')",
            "",
            "def get_embedding(text: str) -> list[float]:",
            "    result = genai.embed_content(",
            "        model='models/text-embedding-004',",
            "        content=text,",
            "        task_type='retrieval_document'  # 검색용 임베딩",
            "    )",
            "    return result['embedding']",
            "",
            "def cosine_similarity(a, b) -> float:",
            "    a, b = np.array(a), np.array(b)",
            "    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b))",
            "",
            "emb1 = get_embedding('파이썬 프로그래밍 언어')",
            "emb2 = get_embedding('파이썬 코딩 튜토리얼')",
            "emb3 = get_embedding('자동차 엔진 수리')",
            "print(cosine_similarity(emb1, emb2))  # → ~0.93",
            "print(cosine_similarity(emb1, emb3))  # → ~0.55",
        ],
        "task_type='retrieval_document'는 검색 인덱싱용, 'retrieval_query'는 쿼리용."
    )

    add_content_slide(prs, "벡터 데이터베이스 (Vector DB)", [
        "## 왜 벡터 DB가 필요한가?",
        "수백만 개 문서 임베딩을 메모리에 올려두기 불가능",
        "빠른 유사도 검색 (ANN: 근사 최근접 이웃 알고리즘) 필요",
        "---",
        "## 주요 벡터 DB 비교",
        "  ChromaDB: 로컬 설치, 쉬운 API, 소규모 프로젝트 ★",
        "  FAISS: Facebook 오픈소스, 매우 빠름, 영구 저장 없음",
        "  Pinecone: 클라우드 관리형, 대규모 상용 서비스",
        "  Weaviate / Qdrant: 자체 호스팅 가능, 고급 기능",
        "---",
        "## ChromaDB 선택 이유",
        "  pip install chromadb 한 줄 설치",
        "  영구 저장 지원 (로컬 폴더)",
        "  Python API가 직관적",
    ])

    add_code_slide(prs, "ChromaDB 기본 사용법",
        [
            "import chromadb",
            "from chromadb.utils import embedding_functions",
            "",
            "# 클라이언트 생성 (로컬 저장)",
            "client = chromadb.PersistentClient(path='./chroma_db')",
            "",
            "# 컬렉션 생성 (테이블과 유사)",
            "collection = client.get_or_create_collection(",
            "    name='my_docs',",
            "    metadata={'hnsw:space': 'cosine'}  # 코사인 유사도",
            ")",
            "",
            "# 문서 추가",
            "collection.add(",
            "    documents=['파이썬은 쉬운 언어입니다', '자바는 강타입 언어입니다'],",
            "    ids=['doc1', 'doc2']",
            ")",
            "",
            "# 유사 문서 검색",
            "results = collection.query(",
            "    query_texts=['프로그래밍 언어 추천'],",
            "    n_results=2",
            ")",
            "print(results['documents'])",
        ],
        "PersistentClient를 사용하면 프로그램 종료 후에도 데이터가 유지됩니다."
    )

    add_summary_slide(prs, 5,
        [
            "임베딩 = 텍스트를 수치 벡터로 변환해 의미 비교 가능",
            "코사인 유사도로 두 텍스트의 의미 거리 측정",
            "벡터 DB = 대량 임베딩의 빠른 유사도 검색 인프라",
            "ChromaDB: pip install 한 줄, 로컬 영구 저장 지원",
            "이것이 RAG의 검색(Retrieval) 핵심 기술",
        ],
        "6주차 예고:\nPDF·텍스트 파싱\n청킹(Chunking) 전략\n완전한 RAG 파이프라인 구현"
    )

    return prs

def make_week06():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 6,
        "PDF·텍스트 파싱 및\n지식 검색 시스템 구현",
        "RAG 파이프라인 완성 — 내 문서로 답하는 AI",
        "Phase 1 — RAG 기초")

    add_agenda_slide(prs, [
        "RAG 파이프라인 전체 구조",
        "문서 로딩 (PDF·TXT·DOCX)",
        "청킹(Chunking) 전략",
        "임베딩 및 벡터 DB 저장",
        "검색 및 답변 생성",
        "실습: 강의 자료 Q&A 봇",
        "RAG 품질 개선 기법",
        "과제 안내",
    ])

    add_content_slide(prs, "RAG 파이프라인 전체 구조", [
        "## RAG = Retrieval-Augmented Generation",
        "LLM의 '지식 한계'를 외부 문서로 보완하는 기법",
        "---",
        "## 2단계 프로세스",
        "  [인덱싱 단계] 문서 로딩 → 청킹 → 임베딩 → 벡터DB 저장",
        "  [검색 단계]  질문 임베딩 → 유사 청크 검색 → LLM 답변 생성",
        "---",
        "## 인덱싱 (오프라인, 1회)",
        "1. 문서 파일 읽기 (PDF, TXT, DOCX …)",
        "2. 적절한 크기로 분할 (청킹)",
        "3. 각 청크를 임베딩 벡터로 변환",
        "4. 벡터 DB에 저장",
        "---",
        "## 검색 (온라인, 매 질문)",
        "1. 사용자 질문 → 임베딩",
        "2. 벡터 DB에서 Top-K 유사 청크 검색",
        "3. [질문 + 검색된 청크]를 LLM에 전달 → 답변",
    ])

    add_content_slide(prs, "문서 로딩 — 다양한 형식 지원", [
        "## 필요 라이브러리",
        "  pip install pypdf2 python-docx",
        "---",
        "## PDF 파싱",
        "  import PyPDF2",
        "  reader = PyPDF2.PdfReader('file.pdf')",
        "  text = ' '.join([p.extract_text() for p in reader.pages])",
        "---",
        "## TXT / Markdown",
        "  with open('file.txt', encoding='utf-8') as f: text = f.read()",
        "---",
        "## DOCX (Word)",
        "  from docx import Document",
        "  doc = Document('file.docx')",
        "  text = '\\n'.join([p.text for p in doc.paragraphs])",
        "---",
        "## 주의사항",
        "  스캔 PDF → OCR 도구 필요 (pdf2image + pytesseract)",
        "  표·이미지 내 텍스트는 별도 처리 필요",
    ])

    add_content_slide(prs, "청킹(Chunking) 전략", [
        "## 왜 청킹이 중요한가?",
        "너무 큰 청크: 관련 없는 정보 포함 → 노이즈",
        "너무 작은 청크: 맥락 손실 → 부정확한 검색",
        "---",
        "## 전략 1: 고정 크기 청킹",
        "  chunk_size=500 토큰, overlap=50 토큰",
        "  구현이 단순, 문장 중간 잘림 가능",
        "---",
        "## 전략 2: 문장/단락 기준",
        "  마침표·줄바꿈 기준으로 분할 → 자연스러운 경계",
        "---",
        "## 전략 3: 의미 기반 청킹 (Semantic)",
        "  연속 문장의 임베딩 유사도가 낮아지는 지점에서 분할",
        "  가장 정확하지만 구현이 복잡",
        "---",
        "## 권장: chunk_size=400~800, overlap=50~100",
    ])

    add_code_slide(prs, "완전한 RAG 파이프라인 구현",
        [
            "import chromadb, PyPDF2, google.generativeai as genai",
            "",
            "# ── 인덱싱 ────────────────────────────────",
            "def build_index(pdf_path: str):",
            "    reader = PyPDF2.PdfReader(pdf_path)",
            "    full_text = ' '.join(p.extract_text() for p in reader.pages)",
            "    # 500자 단위 청킹 (overlap 50)",
            "    chunks = [full_text[i:i+500] for i in range(0, len(full_text), 450)]",
            "    client = chromadb.PersistentClient('./rag_db')",
            "    col = client.get_or_create_collection('docs')",
            "    col.add(documents=chunks, ids=[f'c{i}' for i in range(len(chunks))])",
            "    return col",
            "",
            "# ── 검색 + 답변 ───────────────────────────",
            "def ask(col, model, question: str) -> str:",
            "    results = col.query(query_texts=[question], n_results=3)",
            "    context = '\\n---\\n'.join(results['documents'][0])",
            "    prompt = f'''다음 문서를 참고해 질문에 답하세요.",
            "문서:\\n{context}\\n\\n질문: {question}'''",
            "    return model.generate_content(prompt).text",
        ],
        "RAG의 핵심 로직 — 인덱싱(1회)과 검색+답변(매 질문)을 분리합니다."
    )

    add_content_slide(prs, "RAG 품질 개선 기법", [
        "## 검색 품질 개선",
        "  Hybrid Search: 키워드 검색 + 벡터 검색 결합",
        "  MMR (Maximal Marginal Relevance): 다양성 높은 청크 선택",
        "  Re-ranking: 검색 결과 재정렬 (Cross-Encoder 모델)",
        "---",
        "## 프롬프트 품질 개선",
        "  '문서에 없으면 모른다고 답하세요' — 환각 방지",
        "  출처(어떤 문서 몇 번째 청크)를 함께 반환",
        "---",
        "## 청킹 품질 개선",
        "  Parent-Child 청킹: 큰 청크로 검색, 작은 청크로 답변",
        "  메타데이터 필터: 날짜·카테고리로 사전 필터링",
        "---",
        "## 평가",
        "  RAGAS 라이브러리 — Faithfulness, Answer Relevancy 측정",
    ])

    add_summary_slide(prs, 6,
        [
            "RAG = 문서 인덱싱(1회) + 검색+생성(매 질문)",
            "청킹: 400~800자, 50자 오버랩이 일반적인 시작점",
            "PyPDF2로 PDF 텍스트 추출 → ChromaDB에 저장",
            "프롬프트에 '모르면 모른다고 답하라' 추가로 환각 방지",
            "Phase 1 RAG 완료! 다음은 외부 API 연동(Function Calling)",
        ],
        "7주차 예고:\nFunction Calling\n날씨·검색 API 연동\n도구 사용 에이전트"
    )

    return prs

def make_week07():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 7,
        "Function Calling\n(도구 연동)",
        "LLM에게 외부 세계를 연결하다",
        "Phase 1 — Function Calling")

    add_agenda_slide(prs, [
        "Function Calling이란?",
        "동작 원리 상세",
        "도구 정의 (Tool Schema)",
        "날씨 API 연동 실습",
        "검색 API 연동",
        "다중 도구(Multi-tool) 에이전트",
        "에러 처리 및 안전장치",
        "Phase 1 총 마무리",
    ])

    add_content_slide(prs, "Function Calling이란?", [
        "## 개념",
        "LLM이 직접 함수를 실행하는 것이 아니라,",
        "'이 함수를 실행해달라'는 구조화된 요청을 반환 → 코드가 실행 후 결과 전달",
        "---",
        "## 왜 필요한가?",
        "  LLM은 실시간 정보 없음 (학습 데이터 이후 정보 모름)",
        "  계산, 검색, DB 조회 등 외부 작업 불가",
        "---",
        "## 동작 흐름",
        "1. 사용자: '오늘 서울 날씨 알려줘'",
        "2. LLM: get_weather(city='서울') 호출 요청 반환",
        "3. 코드: 실제 날씨 API 호출",
        "4. 결과를 LLM에 전달 → 자연어 답변 생성",
        "---",
        "## LLM의 역할",
        "언제 어떤 함수를 호출할지 결정 (Router 역할)",
    ])

    add_code_slide(prs, "도구 정의 및 Function Calling 구현",
        [
            "import google.generativeai as genai",
            "import requests",
            "",
            "# 1. 도구 함수 정의",
            "def get_weather(city: str) -> dict:",
            "    '''지정한 도시의 현재 날씨를 반환합니다.'''",
            "    # 실제 환경: OpenWeatherMap API 호출",
            "    return {'city': city, 'temp': 22, 'condition': '맑음'}",
            "",
            "# 2. 모델에 도구 등록",
            "model = genai.GenerativeModel(",
            "    'gemini-1.5-flash',",
            "    tools=[get_weather]  # 함수 직접 전달 (자동 스키마 추출)",
            ")",
            "",
            "# 3. 대화 시작",
            "chat = model.start_chat(enable_automatic_function_calling=True)",
            "response = chat.send_message('서울 날씨 어때?')",
            "print(response.text)",
            "# → '현재 서울은 22°C이며 맑은 날씨입니다.'",
        ],
        "enable_automatic_function_calling=True로 함수 실행을 자동화할 수 있습니다."
    )

    add_content_slide(prs, "수동 Function Calling — 완전한 제어", [
        "## 자동 vs 수동",
        "  자동: SDK가 함수 실행까지 처리 (간편하지만 제어 어려움)",
        "  수동: 코드가 직접 함수 호출 결정 (권장 — 에러 처리 용이)",
        "---",
        "## 수동 처리 흐름",
        "1. response = chat.send_message(user_input)",
        "2. if response.candidates[0].content.parts[0].function_call:",
        "3.     fn_call = response.candidates[0].content.parts[0].function_call",
        "4.     result = execute_function(fn_call.name, fn_call.args)",
        "5.     response = chat.send_message(FunctionResponse(result))",
        "---",
        "## 장점",
        "  함수 실행 전 로그 기록",
        "  실패 시 폴백 처리",
        "  보안 검증 (허용된 함수만 실행)",
    ])

    add_code_slide(prs, "다중 도구 에이전트 구현",
        [
            "def search_web(query: str) -> str:",
            "    '''웹에서 최신 정보를 검색합니다.'''",
            "    return f'검색 결과: {query}에 대한 최신 정보...'",
            "",
            "def calculate(expression: str) -> float:",
            "    '''수학 계산을 수행합니다. 예: 2+3*4'''",
            "    return eval(expression)  # 실제: 안전한 파서 사용",
            "",
            "def get_weather(city: str) -> dict:",
            "    '''날씨 정보를 반환합니다.'''",
            "    return {'city': city, 'temp': 20, 'condition': '흐림'}",
            "",
            "# 세 가지 도구 모두 등록",
            "model = genai.GenerativeModel(",
            "    'gemini-1.5-flash',",
            "    tools=[search_web, calculate, get_weather],",
            "    system_instruction='사용자 요청에 적합한 도구를 선택해 답하세요.'",
            ")",
            "",
            "chat = model.start_chat(enable_automatic_function_calling=True)",
            "# '내일 부산 날씨와 지금 환율 알려줘' → 두 도구 동시 호출 가능",
        ],
        "여러 도구를 등록하면 LLM이 상황에 따라 적절한 도구를 선택합니다."
    )

    add_summary_slide(prs, 7,
        [
            "Function Calling = LLM이 '함수 호출 요청'을 반환, 코드가 실행",
            "도구 함수에 docstring 필수 — LLM이 도구를 이해하는 근거",
            "enable_automatic_function_calling으로 자동 실행 가능",
            "수동 처리로 로그·보안·에러 처리를 세밀하게 제어",
            "Phase 1 완료! 1~7주 배운 모든 기술이 에이전트의 재료",
        ],
        "8주차 예고:\n에이전트란?\nReAct 패턴\nLangGraph 기초\n프로젝트 시작!"
    )

    return prs

def make_week08():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 8,
        "에이전트 입문 &\n프로젝트 시작",
        "ReAct 패턴 + LangGraph 기초 + 프로젝트 기획",
        "Phase 2 — 프로젝트 킥오프 (프로젝트 1주차)")

    add_agenda_slide(prs, [
        "Phase 2 오리엔테이션",
        "에이전트란? 정의와 구성 요소",
        "ReAct 패턴 — Think·Act·Observe",
        "LangGraph 소개 및 설치",
        "첫 번째 LangGraph 그래프",
        "프로젝트 주제 브레인스토밍",
        "기획서 작성 가이드",
        "과제: 프로젝트 기획서 제출",
    ])

    add_content_slide(prs, "에이전트(Agent)란?", [
        "## 기존 LLM vs 에이전트",
        "  기존 LLM: 질문 → 답변 (1회성, 수동적)",
        "  에이전트: 목표 → 계획 → 행동 → 관찰 → 반복 (자율적)",
        "---",
        "## 에이전트의 4가지 구성 요소",
        "  두뇌(Brain): LLM — 추론 및 결정",
        "  도구(Tools): 검색·계산·DB·API 등",
        "  메모리(Memory): 단기(대화) + 장기(벡터DB)",
        "  계획(Planning): 목표를 하위 작업으로 분해",
        "---",
        "## 에이전트 활용 사례",
        "  자동 코드 리뷰 에이전트",
        "  고객 지원 에이전트 (RAG + API 연동)",
        "  연구 보조 에이전트 (검색 + 요약 + 정리)",
    ])

    add_content_slide(prs, "ReAct 패턴 — 에이전트의 사고 방식", [
        "## ReAct = Reasoning + Acting",
        "논문: 'ReAct: Synergizing Reasoning and Acting in Language Models' (2022)",
        "---",
        "## 루프 구조",
        "  Thought(생각): '날씨를 알려면 날씨 API를 써야 해'",
        "  Action(행동): get_weather(city='서울') 호출",
        "  Observation(관찰): '22°C, 맑음'",
        "  → 목표 달성 시 최종 답변 반환",
        "---",
        "## 실제 동작 예시",
        "  목표: '서울과 부산 중 어디가 더 덥고 이유는?'",
        "  1) Thought: 두 도시 날씨 정보 필요",
        "  2) Action: get_weather('서울'), get_weather('부산')",
        "  3) Observation: 서울 28°C, 부산 31°C",
        "  4) Thought: 부산이 더 더움. 해양성 기후 + 복사열 원인",
        "  5) Final Answer: '부산이 31°C로 더 덥습니다. ...'",
    ])

    add_content_slide(prs, "LangGraph 소개", [
        "## LangGraph란?",
        "LLM 에이전트의 복잡한 흐름을 그래프(DAG) 형태로 설계하는 프레임워크",
        "  개발사: LangChain (langgraph.com)",
        "  설치: pip install langgraph",
        "---",
        "## 핵심 개념",
        "  상태(State): 그래프 전체가 공유하는 데이터 (TypedDict)",
        "  노드(Node): 상태를 변환하는 함수",
        "  엣지(Edge): 노드 간 흐름 (조건부 엣지 지원)",
        "---",
        "## LangGraph vs 단순 루프",
        "  단순 루프: 직선 흐름, 디버깅 어려움",
        "  LangGraph: 시각화, 조건 분기, 상태 추적, 스트리밍 지원",
        "---",
        "## 설치",
        "  pip install langgraph langchain-google-genai",
    ])

    add_code_slide(prs, "첫 번째 LangGraph 에이전트",
        [
            "from typing import TypedDict, Annotated",
            "from langgraph.graph import StateGraph, END",
            "from langchain_google_genai import ChatGoogleGenerativeAI",
            "from langchain_core.messages import HumanMessage",
            "",
            "class AgentState(TypedDict):",
            "    messages: list",
            "    next_action: str",
            "",
            "llm = ChatGoogleGenerativeAI(model='gemini-1.5-flash')",
            "",
            "def think_node(state: AgentState) -> AgentState:",
            "    response = llm.invoke(state['messages'])",
            "    return {'messages': state['messages'] + [response]}",
            "",
            "# 그래프 구성",
            "graph = StateGraph(AgentState)",
            "graph.add_node('think', think_node)",
            "graph.set_entry_point('think')",
            "graph.add_edge('think', END)",
            "app = graph.compile()",
            "",
            "result = app.invoke({'messages': [HumanMessage('안녕!')], 'next_action': ''})",
            "print(result['messages'][-1].content)",
        ],
        "State → Node → Edge의 3요소로 LangGraph 에이전트를 구성합니다."
    )

    add_content_slide(prs, "프로젝트 기획 가이드", [
        "## 프로젝트 요구사항",
        "  RAG + 에이전트 기술을 결합한 서비스",
        "  Streamlit으로 웹 UI 제공",
        "  실제 배포 (Streamlit Cloud / HuggingFace)",
        "---",
        "## 좋은 프로젝트 아이디어 조건",
        "  내가 실제로 쓰고 싶은 서비스",
        "  명확한 데이터셋 (PDF, 텍스트 등)",
        "  기술적으로 도전적이지만 달성 가능",
        "---",
        "## 아이디어 예시",
        "  대학 강의 자료 Q&A 챗봇",
        "  특정 법령/판례 검색 도우미",
        "  게임 공략집 챗봇",
        "  특정 분야 논문 요약 서비스",
        "---",
        "## 기획서 포함 항목",
        "  서비스명, 목표 사용자, 핵심 기능, 데이터셋, 기술 스택",
    ])

    add_summary_slide(prs, 8,
        [
            "에이전트 = 두뇌(LLM) + 도구 + 메모리 + 계획의 결합체",
            "ReAct = Thought → Action → Observation 반복 루프",
            "LangGraph = 에이전트 흐름을 그래프로 설계하는 프레임워크",
            "State(공유 데이터) → Node(함수) → Edge(흐름) 3요소",
            "이번 주 과제: 프로젝트 기획서 작성 (데이터셋 포함)",
        ],
        "9주차 예고:\nLangGraph 심화\n복잡한 조건 분기\nRAG+도구 통합 설계"
    )

    return prs

def make_week09():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 9,
        "에이전트 상태 관리 및\n워크플로우",
        "LangGraph 심화 — 조건 분기·루프·도구 통합",
        "Phase 2 — 프로젝트 2주차")

    add_agenda_slide(prs, [
        "지난 주 기획서 피드백",
        "LangGraph 상태 심화",
        "조건부 엣지 (라우팅)",
        "루프와 종료 조건",
        "도구 노드 통합",
        "RAG 노드 구현",
        "실습: 프로젝트 에이전트 로직 구현",
        "과제: 도구 사용 논리 구조 완성",
    ])

    add_content_slide(prs, "LangGraph 상태(State) 심화", [
        "## State = 에이전트의 '기억'",
        "모든 노드가 공유하는 TypedDict — 그래프 실행 내내 유지",
        "---",
        "## Annotated + reducer로 리스트 누적",
        "  messages: Annotated[list, operator.add]",
        "  → 각 노드가 반환한 messages가 기존 리스트에 추가(override가 아님)",
        "---",
        "## 커스텀 상태 필드",
        "  current_tool: str — 현재 사용 중인 도구",
        "  retrieved_docs: list — 검색된 문서 청크",
        "  user_intent: str — 사용자 의도 분류 결과",
        "  iteration_count: int — 루프 횟수 추적",
        "---",
        "## 상태 설계 팁",
        "  필요한 정보만 최소한으로 유지",
        "  노드 간 의존성을 상태로 표현",
    ])

    add_code_slide(prs, "조건부 엣지로 라우팅 구현",
        [
            "from langgraph.graph import StateGraph, END",
            "from typing import Literal",
            "",
            "class AgentState(TypedDict):",
            "    messages: list",
            "    user_intent: str  # 'rag' | 'search' | 'chat'",
            "",
            "def classify_intent(state: AgentState) -> AgentState:",
            "    '''사용자 의도를 분류하는 노드'''",
            "    last_msg = state['messages'][-1].content",
            "    if '문서' in last_msg or '자료' in last_msg:",
            "        return {**state, 'user_intent': 'rag'}",
            "    elif '검색' in last_msg or '최신' in last_msg:",
            "        return {**state, 'user_intent': 'search'}",
            "    else:",
            "        return {**state, 'user_intent': 'chat'}",
            "",
            "def route(state: AgentState) -> Literal['rag_node', 'search_node', 'chat_node']:",
            "    return f\"{state['user_intent']}_node\"",
            "",
            "graph = StateGraph(AgentState)",
            "graph.add_conditional_edges('classify', route, {",
            "    'rag_node': 'rag_node',",
            "    'search_node': 'search_node',",
            "    'chat_node': 'chat_node'",
            "})",
        ],
        "add_conditional_edges로 상태에 따라 다른 노드로 분기합니다."
    )

    add_code_slide(prs, "RAG 노드 + 도구 노드 통합",
        [
            "import chromadb",
            "",
            "chroma_client = chromadb.PersistentClient('./rag_db')",
            "collection = chroma_client.get_collection('docs')",
            "",
            "def rag_node(state: AgentState) -> AgentState:",
            "    '''벡터 DB에서 관련 문서를 검색하고 답변 생성'''",
            "    query = state['messages'][-1].content",
            "    # 1. 검색",
            "    results = collection.query(query_texts=[query], n_results=3)",
            "    context = '\\n'.join(results['documents'][0])",
            "    # 2. 답변 생성",
            "    prompt = f'다음 문서 기반으로 답하세요:\\n{context}\\n\\n질문: {query}'",
            "    answer = llm.invoke([HumanMessage(prompt)])",
            "    return {**state, 'messages': state['messages'] + [answer]}",
            "",
            "def search_node(state: AgentState) -> AgentState:",
            "    '''웹 검색 도구를 사용해 최신 정보 제공'''",
            "    # 실제: Tavily / DuckDuckGo API 연동",
            "    query = state['messages'][-1].content",
            "    search_result = f'{query}에 대한 최신 검색 결과...'",
            "    answer = llm.invoke([HumanMessage(f'검색결과: {search_result}')])",
            "    return {**state, 'messages': state['messages'] + [answer]}",
        ],
        "RAG 노드와 검색 노드를 분리하면 각 경로를 독립적으로 개선할 수 있습니다."
    )

    add_content_slide(prs, "루프와 종료 조건", [
        "## ReAct 루프 구현",
        "에이전트가 목표를 달성할 때까지 반복 실행",
        "  think → act → observe → think → ... → 종료",
        "---",
        "## 종료 조건",
        "  LLM이 'FINAL ANSWER' 키워드 반환",
        "  최대 반복 횟수 초과 (안전장치, 보통 10~20회)",
        "  더 이상 호출할 도구 없음",
        "---",
        "## 무한 루프 방지",
        "  state['iteration_count'] >= MAX_ITERATIONS → END",
        "---",
        "## LangGraph의 루프 구현",
        "  조건부 엣지: should_continue() → 'continue' or 'end'",
        "  'end' → END 노드로 이동",
        "---",
        "## 실습 목표",
        "  프로젝트 에이전트의 전체 그래프 설계 및 구현",
    ])

    add_summary_slide(prs, 9,
        [
            "State는 TypedDict — Annotated + reducer로 리스트 누적",
            "add_conditional_edges로 상태 기반 동적 라우팅",
            "RAG 노드와 검색 노드를 분리해 독립적으로 개선",
            "루프 종료 조건 필수 — MAX_ITERATIONS로 안전장치",
            "이번 주 과제: 프로젝트 에이전트 도구 로직 완성",
        ],
        "10주차 예고:\nStreamlit 웹 UI\n채팅 인터페이스 구축\n세션 상태 관리"
    )

    return prs

def make_week10():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 10,
        "Streamlit 웹 UI 구현",
        "채팅 인터페이스 + 세션 상태 + 에이전트 연결",
        "Phase 2 — 프로젝트 3주차")

    add_agenda_slide(prs, [
        "Streamlit이란?",
        "기본 UI 컴포넌트",
        "st.chat_message & st.chat_input",
        "st.session_state 완전 정복",
        "파일 업로드 처리",
        "에이전트 백엔드 연결",
        "실습: 챗봇 UI 완성",
        "과제: 프론트-백엔드 통합",
    ])

    add_content_slide(prs, "Streamlit이란?", [
        "## 개요",
        "Python 코드만으로 웹 앱을 만드는 프레임워크",
        "HTML/CSS/JS 없이 데이터 앱·챗봇 UI 구현 가능",
        "  설치: pip install streamlit",
        "  실행: streamlit run app.py",
        "---",
        "## 특징",
        "  코드 변경 시 자동 새로고침 (Hot Reload)",
        "  Streamlit Cloud로 1-click 무료 배포",
        "  chat_message, chat_input 등 챗봇 전용 컴포넌트 내장",
        "---",
        "## 실행 모델",
        "  사용자 상호작용마다 스크립트 전체 재실행",
        "  → st.session_state로 상태 유지 필수",
        "---",
        "## 주요 컴포넌트",
        "  st.title, st.write, st.button, st.text_input",
        "  st.chat_message, st.chat_input, st.sidebar",
    ])

    add_code_slide(prs, "기본 챗봇 UI 구조",
        [
            "import streamlit as st",
            "",
            "st.title('🤖 AI 챗봇')",
            "",
            "# 대화 이력 초기화",
            "if 'messages' not in st.session_state:",
            "    st.session_state.messages = []",
            "",
            "# 이전 메시지 표시",
            "for msg in st.session_state.messages:",
            "    with st.chat_message(msg['role']):",
            "        st.write(msg['content'])",
            "",
            "# 사용자 입력 처리",
            "if prompt := st.chat_input('메시지를 입력하세요...'):",
            "    # 사용자 메시지 추가 및 표시",
            "    st.session_state.messages.append({'role': 'user', 'content': prompt})",
            "    with st.chat_message('user'):",
            "        st.write(prompt)",
            "",
            "    # AI 응답 생성 및 표시",
            "    with st.chat_message('assistant'):",
            "        response = '안녕하세요! 무엇을 도와드릴까요?'  # 임시",
            "        st.write(response)",
            "    st.session_state.messages.append({'role': 'assistant', 'content': response})",
        ],
        "if 'key' not in st.session_state: 패턴으로 상태를 초기화합니다."
    )

    add_code_slide(prs, "LangGraph 에이전트 연결",
        [
            "import streamlit as st",
            "from agent import create_agent  # 에이전트 생성 함수",
            "",
            "# 에이전트 초기화 (앱 로드 시 1회만)",
            "@st.cache_resource",
            "def load_agent():",
            "    return create_agent()",
            "",
            "agent = load_agent()",
            "",
            "if prompt := st.chat_input('질문하세요...'):",
            "    # ... 사용자 메시지 표시 ...",
            "",
            "    with st.chat_message('assistant'):",
            "        with st.spinner('생각 중...'):",
            "            # 에이전트 실행",
            "            result = agent.invoke({",
            "                'messages': [{'role': 'user', 'content': prompt}]",
            "            })",
            "            answer = result['messages'][-1].content",
            "        st.write(answer)",
            "",
            "    st.session_state.messages.append({",
            "        'role': 'assistant', 'content': answer",
            "    })",
        ],
        "@st.cache_resource로 에이전트를 앱 전체에서 1번만 초기화합니다."
    )

    add_content_slide(prs, "파일 업로드 및 사이드바 구성", [
        "## 파일 업로드",
        "  uploaded = st.sidebar.file_uploader('PDF 업로드', type=['pdf'])",
        "  if uploaded: process_pdf(uploaded)",
        "---",
        "## 사이드바 구성",
        "  모델 선택, 온도 조절, 업로드 등 설정 UI",
        "  with st.sidebar: → 사이드바 컴포넌트 배치",
        "---",
        "## 스트리밍 답변 표시",
        "  with st.chat_message('assistant'):",
        "    response_placeholder = st.empty()",
        "    full_response = ''",
        "    for chunk in agent.stream(input):  ",
        "        full_response += chunk",
        "        response_placeholder.write(full_response + '▌')",
        "---",
        "## 유용한 컴포넌트",
        "  st.spinner() — 로딩 표시",
        "  st.toast() — 알림 메시지",
        "  st.expander() — 접기/펼치기",
    ])

    add_summary_slide(prs, 10,
        [
            "Streamlit: Python만으로 웹 챗봇 UI 구현 가능",
            "st.session_state로 새로고침 후에도 대화 이력 유지",
            "@st.cache_resource로 에이전트 초기화 1회만 실행",
            "st.chat_message + st.chat_input = 챗봇 UI의 핵심",
            "이번 주 과제: 에이전트 백엔드와 Streamlit UI 연결",
        ],
        "11주차 예고:\n전체 파이프라인 통합\n디버깅 기법\n품질 개선"
    )

    return prs

def make_week11():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 11,
        "시스템 통합 및\n디버깅",
        "전체 파이프라인 통합 테스트 + 버그 수정",
        "Phase 2 — 프로젝트 4주차")

    add_agenda_slide(prs, [
        "전체 아키텍처 리뷰",
        "공통 버그 패턴 분석",
        "RAG 답변 품질 개선",
        "에이전트 루프 예외 처리",
        "LangGraph 디버깅 도구",
        "로깅 및 모니터링",
        "실습: 개인 프로젝트 통합 테스트",
        "과제: 버그 수정 리포트",
    ])

    add_content_slide(prs, "공통 버그 패턴 Top 5", [
        "## 1. API Key 관련",
        "  환경변수 로드 실패 → load_dotenv() 위치 확인",
        "---",
        "## 2. 벡터 DB 경로 오류",
        "  PersistentClient('./db') — 상대 경로 주의 (CWD 기준)",
        "---",
        "## 3. 세션 상태 초기화 누락",
        "  if 'key' not in st.session_state: 초기화 코드 빠짐",
        "---",
        "## 4. LangGraph 상태 불변성 위반",
        "  state['messages'].append() 대신 return {**state, 'messages': [...]}",
        "---",
        "## 5. 스트리밍 중 .text 접근",
        "  stream=True인 상태에서 response.text 사용 → AttributeError",
        "  → chunk.text 사용 또는 response.resolve() 후 접근",
    ])

    add_content_slide(prs, "RAG 답변 품질 개선", [
        "## 문제: '모르겠습니다' 남발",
        "  원인: 청크가 너무 작거나 임베딩 모델 미스매치",
        "  해결: 청크 크기 늘리기 / 오버랩 증가",
        "---",
        "## 문제: 관련 없는 청크 검색",
        "  원인: n_results 너무 크거나 유사도 임계값 없음",
        "  해결: n_results 줄이기 / where 필터 추가",
        "---",
        "## 문제: 환각(Hallucination)",
        "  모델이 문서에 없는 내용을 만들어냄",
        "  해결: '문서에 없으면 모른다고 답하세요' 프롬프트 추가",
        "---",
        "## 문제: 느린 검색 속도",
        "  원인: 임베딩을 검색마다 재생성",
        "  해결: @st.cache_data로 임베딩 캐시",
    ])

    add_code_slide(prs, "에이전트 예외 처리 패턴",
        [
            "import traceback",
            "from typing import Optional",
            "",
            "def safe_tool_call(func, *args, **kwargs) -> tuple[bool, any]:",
            "    '''도구 호출 안전 래퍼'''",
            "    try:",
            "        result = func(*args, **kwargs)",
            "        return True, result",
            "    except Exception as e:",
            "        error_msg = f'오류 발생: {type(e).__name__}: {str(e)}'",
            "        print(traceback.format_exc())  # 개발자용 로그",
            "        return False, error_msg",
            "",
            "def tool_node(state: AgentState) -> AgentState:",
            "    tool_name = state.get('current_tool')",
            "    tool_args = state.get('tool_args', {})",
            "    ",
            "    success, result = safe_tool_call(TOOLS[tool_name], **tool_args)",
            "    ",
            "    if not success:",
            "        # 실패 시 에이전트에게 오류 알림",
            "        return {**state, 'tool_result': result, 'error': True}",
            "    ",
            "    return {**state, 'tool_result': result, 'error': False}",
        ],
        "도구 실행 실패를 우아하게 처리해 에이전트가 대안을 찾도록 합니다."
    )

    add_content_slide(prs, "LangGraph 디버깅 도구", [
        "## 그래프 시각화",
        "  from IPython.display import Image",
        "  Image(app.get_graph().draw_mermaid_png())",
        "---",
        "## 스트림으로 상태 추적",
        "  for event in app.stream(input, stream_mode='values'):",
        "      print(event)  # 각 노드 실행 후 상태 출력",
        "---",
        "## 체크포인트 (Checkpointing)",
        "  from langgraph.checkpoint.memory import MemorySaver",
        "  app = graph.compile(checkpointer=MemorySaver())",
        "  → 중간 상태 저장 및 특정 지점에서 재실행 가능",
        "---",
        "## print 디버깅",
        "  각 노드 진입·퇴출 시 상태 주요 필드 출력",
        "  f'[{node_name}] intent={state[\"user_intent\"]}'",
    ])

    add_summary_slide(prs, 11,
        [
            "상태 불변성: LangGraph 노드는 항상 새 dict 반환",
            "safe_tool_call 래퍼로 도구 실패를 우아하게 처리",
            "RAG 품질: 청크 크기·오버랩·n_results 조정으로 개선",
            "app.stream(stream_mode='values')로 각 노드 상태 추적",
            "이번 주 과제: 전체 파이프라인 통합 테스트 및 버그 수정 완료",
        ],
        "12주차 예고:\nStreamlit Cloud 배포\nHuggingFace Spaces\nAPI 보안 설정"
    )

    return prs

def make_week12():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 12,
        "클라우드 배포 및\n최종 개발",
        "Streamlit Cloud / HuggingFace Spaces + API 보안",
        "Phase 2 — 프로젝트 5주차")

    add_agenda_slide(prs, [
        "배포 플랫폼 비교",
        "Streamlit Cloud 배포",
        "HuggingFace Spaces 배포",
        "requirements.txt 작성",
        "환경변수 (Secrets) 설정",
        "API 보안 설정",
        "실습: 실제 퍼블릭 URL 배포",
        "Phase 2 마무리",
    ])

    add_two_col_slide(prs, "배포 플랫폼 비교",
        "Streamlit Cloud", [
            "무료 (공개 레포 기준)",
            "GitHub 연동, 1-click 배포",
            "Streamlit 앱에 최적화",
            "자동 재배포 (git push 시)",
            "한계: 리소스 제한, 슬립 모드",
            "",
            "→ 빠른 프로토타입·수업용 ★",
        ],
        "HuggingFace Spaces", [
            "무료 (CPU Basic)",
            "Docker / Gradio / Streamlit 지원",
            "ML 모델 공유에 최적화",
            "커뮤니티 탐색 노출",
            "한계: 빌드 시간 더 김",
            "",
            "→ 모델 쇼케이스·포트폴리오",
        ]
    )

    add_code_slide(prs, "Streamlit Cloud 배포 준비",
        [
            "# 1. requirements.txt 생성",
            "# pip freeze > requirements.txt  (전체 환경)",
            "# 또는 직접 작성 (권장 — 필요한 것만)",
            "",
            "# requirements.txt",
            "streamlit>=1.35.0",
            "google-generativeai>=0.7.0",
            "langchain-google-genai>=1.0.0",
            "langgraph>=0.1.0",
            "chromadb>=0.5.0",
            "pypdf2>=3.0.0",
            "python-dotenv>=1.0.0",
            "",
            "# 2. .streamlit/config.toml (선택)",
            "[theme]",
            "primaryColor = '#89B4FA'",
            "backgroundColor = '#1E1E2E'",
            "secondaryBackgroundColor = '#313244'",
            "textColor = '#FFFFFF'",
            "",
            "# 3. GitHub에 push (단, .env는 절대 포함 금지!)",
            "# .gitignore에 .env, __pycache__, *.pyc, chroma_db/ 추가",
        ],
        "requirements.txt는 버전 범위(>=)로 명시해 의존성 충돌을 예방합니다."
    )

    add_content_slide(prs, "Streamlit Cloud 배포 단계", [
        "## Step 1: GitHub 레포 준비",
        "  main 브랜치에 app.py, requirements.txt 포함",
        "  .gitignore에 .env, chroma_db/ 추가",
        "---",
        "## Step 2: Streamlit Cloud 연결",
        "  share.streamlit.io → 'New app' → GitHub 레포 선택",
        "  메인 파일 경로: app.py",
        "---",
        "## Step 3: Secrets 설정",
        "  App settings → Secrets → TOML 형식으로 입력",
        "  [secrets]",
        "  GOOGLE_API_KEY = 'AIza...'",
        "  코드에서: st.secrets['GOOGLE_API_KEY']",
        "---",
        "## Step 4: 배포 확인",
        "  https://[앱이름].streamlit.app 으로 접속",
        "  git push 시 자동 재배포",
    ])

    add_content_slide(prs, "API 보안 설정", [
        "## 절대 하지 말 것",
        "  API Key를 코드에 하드코딩",
        "  .env 파일을 GitHub에 커밋",
        "  API Key를 로그에 출력",
        "---",
        "## 올바른 방법",
        "  로컬: python-dotenv + .env + .gitignore",
        "  Streamlit Cloud: st.secrets 사용",
        "  HuggingFace: Space Variables/Secrets 사용",
        "---",
        "## 요청 제한 (Rate Limiting)",
        "  Gemini 무료: 분당 15회 → 사용자 세션당 제한 고려",
        "  time.sleep(1) 또는 사용량 카운터로 보호",
        "---",
        "## 입력 검증",
        "  사용자 입력 길이 제한 (max_chars=500)",
        "  악의적 프롬프트 인젝션 방어 (시스템 프롬프트로 역할 고정)",
    ])

    add_summary_slide(prs, 12,
        [
            "Streamlit Cloud: GitHub 연동, Secrets로 API Key 관리",
            ".gitignore에 .env, chroma_db/ 필수 추가",
            "st.secrets['KEY']로 배포 환경 환경변수 접근",
            "API Key는 절대 코드·로그에 노출하지 말 것",
            "Phase 2 완료! 실제 퍼블릭 URL로 서비스 배포 완성",
        ],
        "13주차 예고:\n성능 평가 (RAGAS)\nLangSmith 모니터링\n프롬프트 미세 조정"
    )

    return prs

def make_week13():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 13,
        "성능 평가 및\n모니터링",
        "LangSmith · RAGAS · 프롬프트 미세 조정",
        "Phase 3 — 고도화")

    add_agenda_slide(prs, [
        "LLM 평가가 어려운 이유",
        "RAG 평가 지표 (RAGAS)",
        "LangSmith 소개 및 설정",
        "트레이스 분석 실습",
        "프롬프트 미세 조정",
        "RAG 로직 개선 사례",
        "A/B 테스트 방법론",
        "최종 발표 준비 안내",
    ])

    add_content_slide(prs, "LLM 평가가 어려운 이유", [
        "## 전통적 ML vs LLM 평가",
        "  전통: 정확도(Accuracy) 하나로 측정 가능",
        "  LLM: 정답이 하나가 아님 / 주관적 판단 포함",
        "---",
        "## RAG 특화 문제",
        "  검색이 잘 됐는가? (Retrieval 품질)",
        "  답변이 문서와 일치하는가? (Faithfulness)",
        "  질문에 제대로 답했는가? (Answer Relevancy)",
        "  실제 정답과 일치하는가? (Context Recall)",
        "---",
        "## 평가 접근법",
        "  수동 평가: 사람이 직접 체크 (느리지만 정확)",
        "  LLM-as-Judge: GPT/Claude로 자동 평가 (빠름)",
        "  벤치마크: 공개 데이터셋으로 표준 측정",
    ])

    add_content_slide(prs, "RAGAS — RAG 평가 프레임워크", [
        "## RAGAS란?",
        "RAG 파이프라인을 자동으로 평가하는 오픈소스 라이브러리",
        "  설치: pip install ragas",
        "---",
        "## 핵심 4가지 지표",
        "  Faithfulness (충실도): 답변이 컨텍스트에 기반하는가? (0~1)",
        "  Answer Relevancy (답변 관련성): 질문에 얼마나 관련 있는가?",
        "  Context Precision (컨텍스트 정밀도): 검색된 청크가 유용한가?",
        "  Context Recall (컨텍스트 재현율): 필요한 정보가 검색됐는가?",
        "---",
        "## 평가 데이터셋 구성",
        "  질문(question) + 정답(ground_truth) + AI답변(answer) + 컨텍스트(contexts)",
        "---",
        "## 목표 점수",
        "  Faithfulness > 0.8 / Answer Relevancy > 0.7",
    ])

    add_code_slide(prs, "LangSmith 트레이싱 설정",
        [
            "# 환경변수 설정",
            "import os",
            "os.environ['LANGCHAIN_TRACING_V2'] = 'true'",
            "os.environ['LANGCHAIN_API_KEY'] = 'ls__...'",
            "os.environ['LANGCHAIN_PROJECT'] = 'my-rag-project'",
            "",
            "# 이후 모든 LangChain/LangGraph 실행이 자동으로 추적됨",
            "",
            "# 수동으로 트레이스 추가하고 싶은 경우",
            "from langsmith import traceable",
            "",
            "@traceable(name='RAG 검색')",
            "def retrieve_documents(query: str) -> list:",
            "    results = collection.query(query_texts=[query], n_results=3)",
            "    return results['documents'][0]",
            "",
            "@traceable(name='답변 생성')",
            "def generate_answer(query: str, context: list) -> str:",
            "    prompt = f'컨텍스트: {context}\\n질문: {query}'",
            "    return model.generate_content(prompt).text",
            "",
            "# smith.langchain.com 대시보드에서 실시간 확인",
        ],
        "LANGCHAIN_TRACING_V2=true만 설정하면 모든 실행이 자동으로 추적됩니다."
    )

    add_content_slide(prs, "프롬프트 미세 조정 & A/B 테스트", [
        "## 프롬프트 버전 관리",
        "  LangSmith Hub에 프롬프트 저장·버전 관리",
        "  또는 Python 딕셔너리로 버전 관리",
        "---",
        "## A/B 테스트 방법",
        "1. 10~20개 테스트 질문 수집",
        "2. 프롬프트 버전 A와 B로 각각 답변 생성",
        "3. RAGAS로 자동 평가 또는 직접 채점",
        "4. 더 높은 점수의 버전 선택",
        "---",
        "## 흔한 개선 포인트",
        "  출처 표시 추가: '답변 근거: [문서명 p.X]'",
        "  답변 길이 제어: '3문장 이내로 요약'",
        "  한국어 강제: '반드시 한국어로 답하세요'",
        "---",
        "## RAG 검색 개선",
        "  쿼리 확장: 질문을 여러 형태로 변형 후 검색",
        "  Hybrid Search: 키워드 + 벡터 검색 결합",
    ])

    add_summary_slide(prs, 13,
        [
            "RAGAS 4지표: Faithfulness·Answer Relevancy·Precision·Recall",
            "LangSmith: 환경변수 2개로 자동 트레이싱 시작",
            "A/B 테스트로 프롬프트 개선 효과를 객관적으로 측정",
            "쿼리 확장과 Hybrid Search로 RAG 검색 품질 향상",
            "다음 주 최종 발표 준비: 데모 + 기술 스택 + 배운 점",
        ],
        "14주차:\n최종 프로젝트 데모 발표\n경험 공유\n서로의 서비스 피드백"
    )

    return prs

def make_week14():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 14,
        "최종 프로젝트 데모\n및 발표",
        "라이브 데모 · 기술 스택 · 개발 경험 공유",
        "Phase 3 — 최종 발표")

    add_agenda_slide(prs, [
        "발표 일정 및 규칙",
        "발표 구성 가이드",
        "라이브 데모 준비 체크리스트",
        "기술 발표 팁",
        "상호 피드백 방법",
        "발표 진행 (학생 발표)",
        "종합 피드백",
        "기말고사 안내",
    ])

    add_content_slide(prs, "발표 구성 가이드 (7~10분)", [
        "## 1. 서비스 소개 (1~2분)",
        "  어떤 문제를 해결하는가?",
        "  타겟 사용자는 누구인가?",
        "---",
        "## 2. 라이브 데모 (3~4분)",
        "  실제 배포된 URL로 시연",
        "  핵심 기능 2~3가지 집중 시연",
        "  예상치 못한 결과도 솔직하게 공유",
        "---",
        "## 3. 기술 스택 및 아키텍처 (2~3분)",
        "  사용한 기술 스택 (LLM, RAG, 에이전트, UI)",
        "  전체 데이터 흐름 다이어그램",
        "---",
        "## 4. 도전과 배움 (1~2분)",
        "  개발 중 가장 어려웠던 점",
        "  어떻게 해결했는가?",
        "  다음에 개선하고 싶은 점",
    ])

    add_content_slide(prs, "라이브 데모 체크리스트", [
        "## 배포 전 확인",
        "  배포 URL에서 실제 동작 확인 (로컬 말고)",
        "  API 사용 한도 여유 확인",
        "  발표 당일 인터넷 연결 상태 점검",
        "---",
        "## 데모 준비",
        "  시연용 질문 3~5개 사전 준비",
        "  예상 에러 상황 대비 스크린샷 준비",
        "  UI 설명 없이도 직관적인지 확인",
        "---",
        "## 발표 자료",
        "  아키텍처 다이어그램 (draw.io, Excalidraw 등)",
        "  코드 핵심 부분 스니펫 (최대 1~2개)",
        "---",
        "## 마음가짐",
        "  완벽하지 않아도 괜찮다 — 배운 과정이 더 중요",
        "  질문에 모르면 '모른다'고 솔직하게",
    ])

    add_content_slide(prs, "상호 피드백 방법", [
        "## 피드백 관점",
        "  UX: 사용하기 직관적인가?",
        "  기술: 더 나은 접근법이 있는가?",
        "  창의성: 아이디어가 참신한가?",
        "---",
        "## 좋은 피드백 형식",
        "  '좋았던 점: ...'",
        "  '개선하면 더 좋을 것 같은 점: ...'",
        "  '이 기술을 내 프로젝트에 적용하고 싶은 부분: ...'",
        "---",
        "## 피드백 받는 자세",
        "  방어적으로 반응하지 말 것",
        "  모든 피드백은 서비스 개선을 위한 것",
        "---",
        "## 발표 후 활용",
        "  피드백을 정리해 포트폴리오에 '배운 점' 섹션 추가",
        "  GitHub README 업데이트 (데모 GIF, 기술 스택 명시)",
    ])

    add_content_slide(prs, "포트폴리오로 발전시키기", [
        "## GitHub README 구성",
        "  프로젝트 설명 + 데모 GIF or 스크린샷",
        "  기술 스택 배지 (shields.io)",
        "  설치 및 실행 방법",
        "  아키텍처 다이어그램",
        "---",
        "## 확장 아이디어",
        "  멀티모달 지원 (이미지 질문 처리)",
        "  데이터베이스 연동 (사용자별 대화 저장)",
        "  인증 시스템 추가",
        "  더 많은 도구 연동",
        "---",
        "## 이 수업에서 배운 기술 스택",
        "  LLM API / 프롬프트 엔지니어링",
        "  RAG (ChromaDB + 임베딩)",
        "  에이전트 (LangGraph + ReAct)",
        "  웹 UI (Streamlit) + 클라우드 배포",
    ])

    add_summary_slide(prs, 14,
        [
            "데모는 실제 배포 URL로 — 로컬 환경 시연 금지",
            "기술 스택과 아키텍처 다이어그램을 꼭 준비",
            "어려웠던 점과 해결 과정이 가장 좋은 이야기",
            "상호 피드백으로 모두가 함께 성장",
            "이 프로젝트를 포트폴리오로 발전시키자!",
        ],
        "15주차:\n기말고사\n1~13주 핵심 개념\n이론 중심 시험"
    )

    return prs

def make_week15():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, 15,
        "기말고사",
        "핵심 개념 최종 정리 및 시험 안내",
        "Phase 3 — 기말 평가")

    add_agenda_slide(prs, [
        "기말고사 안내",
        "Part 1: LLM 기초 (1~4주)",
        "Part 2: RAG 파이프라인 (5~6주)",
        "Part 3: 에이전트 (7~9주)",
        "Part 4: 시스템 개발 (10~13주)",
        "예상 문제 유형",
        "마지막 수업 마무리",
        "수고하셨습니다!",
    ])

    add_content_slide(prs, "기말고사 안내", [
        "## 시험 형식",
        "  이론 중심 / 객관식 + 단답형 + 서술형",
        "  총 60~80분 / 100점 만점",
        "---",
        "## 시험 범위",
        "  1~13주차 강의 핵심 개념",
        "  코드 작성 문제 없음 (개념 이해 확인)",
        "---",
        "## 평가 비중",
        "  LLM 기초 (1~4주): 25점",
        "  RAG 파이프라인 (5~6주): 25점",
        "  에이전트 (7~9주): 25점",
        "  시스템 개발·평가 (10~13주): 25점",
        "---",
        "## 준비 방법",
        "  각 주 '핵심 정리' 슬라이드 위주로 복습",
        "  개념의 '왜'와 '어떻게' 연결해서 이해",
    ])

    add_content_slide(prs, "Part 1 — LLM 기초 핵심 정리 (1~4주)", [
        "## LLM 원리",
        "  Transformer Self-Attention / 다음 토큰 예측 / 사전학습+파인튜닝+RLHF",
        "---",
        "## 토큰",
        "  LLM의 최소 처리·비용 단위 / 컨텍스트 윈도우 한계",
        "---",
        "## 프롬프트 엔지니어링",
        "  Zero/Few-shot / Persona (system_instruction) / CoT",
        "  구체적·명확한 지시의 중요성",
        "---",
        "## 구조화된 출력",
        "  response_mime_type='application/json' / Pydantic 검증",
        "---",
        "## 대화 이력",
        "  ChatSession이 히스토리 자동 관리 / 슬라이딩 윈도우로 토큰 절약",
    ])

    add_content_slide(prs, "Part 2 — RAG 파이프라인 핵심 정리 (5~6주)", [
        "## 임베딩",
        "  텍스트 → 수치 벡터 / 의미 유사성 표현 / 코사인 유사도",
        "---",
        "## 벡터 데이터베이스",
        "  대량 임베딩 저장·검색 / ChromaDB (로컬·영구 저장)",
        "  ANN(근사 최근접 이웃) 알고리즘으로 빠른 검색",
        "---",
        "## RAG 파이프라인",
        "  인덱싱: 문서 로딩 → 청킹 → 임베딩 → 벡터DB 저장 (1회)",
        "  검색: 질문 임베딩 → Top-K 검색 → LLM 답변 생성 (매 질문)",
        "---",
        "## 청킹 전략",
        "  400~800자 / 50자 오버랩 / 문장 경계 기준 권장",
        "---",
        "## 품질 개선",
        "  환각 방지 프롬프트 / Hybrid Search / RAGAS 평가",
    ])

    add_content_slide(prs, "Part 3 — 에이전트 핵심 정리 (7~9주)", [
        "## Function Calling",
        "  LLM이 함수 실행 요청 반환 → 코드가 실행 → 결과 전달",
        "  도구 함수 docstring이 LLM의 도구 이해 근거",
        "---",
        "## 에이전트 구성 요소",
        "  두뇌(LLM) + 도구(Tools) + 메모리 + 계획",
        "---",
        "## ReAct 패턴",
        "  Thought → Action → Observation 반복 루프",
        "  최대 반복 횟수 제한 필수 (무한 루프 방지)",
        "---",
        "## LangGraph",
        "  State(TypedDict) → Node(함수) → Edge(흐름)",
        "  add_conditional_edges로 상태 기반 동적 라우팅",
        "  상태 불변성: 노드는 항상 새 dict 반환",
    ])

    add_content_slide(prs, "Part 4 — 시스템 개발 핵심 정리 (10~13주)", [
        "## Streamlit",
        "  매 상호작용마다 스크립트 재실행 → session_state로 상태 유지",
        "  @st.cache_resource로 에이전트 1회만 초기화",
        "---",
        "## 배포",
        "  .gitignore에 .env 필수 / Streamlit Cloud → st.secrets",
        "  requirements.txt에 필요 패키지·버전 명시",
        "---",
        "## API 보안",
        "  API Key 하드코딩 금지 / 환경변수로 관리",
        "---",
        "## 평가 (RAGAS)",
        "  Faithfulness / Answer Relevancy / Context Precision·Recall",
        "---",
        "## LangSmith",
        "  환경변수 2개로 자동 트레이싱 / 노드별 실행 시간·비용 추적",
    ])

    add_content_slide(prs, "예상 문제 유형", [
        "## 객관식 예시",
        "  RAG에서 '청킹'의 목적으로 가장 적절한 것은?",
        "  ReAct 패턴의 올바른 실행 순서는?",
        "  코사인 유사도 값 범위는?",
        "---",
        "## 단답형 예시",
        "  LangGraph에서 State를 정의할 때 사용하는 Python 타입은?",
        "  Streamlit에서 앱 재실행 후에도 상태를 유지하는 객체는?",
        "---",
        "## 서술형 예시",
        "  RAG 파이프라인의 인덱싱 단계를 단계별로 설명하라.",
        "  Function Calling의 동작 원리를 LLM의 역할 중심으로 서술하라.",
        "  슬라이딩 윈도우 히스토리 관리가 필요한 이유를 설명하라.",
        "---",
        "## 단, 코드 작성 문제는 없습니다.",
    ])

    # 마지막 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, YELLOW)
    add_rect(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(4.5), CARD_BG)

    add_textbox(slide, "한 학기 수고하셨습니다!", Inches(0.7), Inches(0.4), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "여러분은 이 수업을 통해", Inches(1.2), Inches(2.0), Inches(11), Inches(0.6),
                font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide,
        "LLM API 제어  ·  프롬프트 엔지니어링  ·  RAG 파이프라인\n"
        "에이전트 개발 (LangGraph)  ·  웹 UI (Streamlit)  ·  클라우드 배포",
        Inches(1.2), Inches(2.7), Inches(11), Inches(1.2),
        font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "...을 직접 만들고 배포하는 경험을 쌓았습니다.",
                Inches(1.2), Inches(4.0), Inches(11), Inches(0.6),
                font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide,
        "이 경험이 여러분의 커리어와 다음 프로젝트에 큰 힘이 되길 바랍니다.",
        Inches(1.2), Inches(5.5), Inches(11), Inches(0.6),
        font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "LLM 기반 AI 서비스 개발 실습", Inches(0.7), Inches(6.8), Inches(12), Inches(0.5),
                font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)
    return prs


# ════════════════════════════════════════════════════════════
#  메인 실행
# ════════════════════════════════════════════════════════════

WEEKS = [
    (make_week01, "Week01_LLM원리및Gemini환경세팅"),
    (make_week02, "Week02_PythonSDK및CLI챗봇구현"),
    (make_week03, "Week03_프롬프트엔지니어링"),
    (make_week04, "Week04_구조화출력및대화이력관리"),
    (make_week05, "Week05_임베딩과벡터데이터베이스"),
    (make_week06, "Week06_PDF파싱및지식검색시스템"),
    (make_week07, "Week07_FunctionCalling도구연동"),
    (make_week08, "Week08_에이전트입문및프로젝트시작"),
    (make_week09, "Week09_에이전트상태관리및워크플로우"),
    (make_week10, "Week10_Streamlit웹UI구현"),
    (make_week11, "Week11_시스템통합및디버깅"),
    (make_week12, "Week12_클라우드배포및최종개발"),
    (make_week13, "Week13_성능평가및모니터링"),
    (make_week14, "Week14_최종프로젝트데모및발표"),
    (make_week15, "Week15_기말고사"),
]

output_dir = "d:/[1]수업자료/PPT"
os.makedirs(output_dir, exist_ok=True)

for i, (make_fn, filename) in enumerate(WEEKS, 1):
    print(f"[{i:02d}/15] 생성 중: {filename}.pptx ...")
    prs = make_fn()
    prs.save(f"{output_dir}/{filename}.pptx")
    print(f"       OK 저장 완료")

print("\n모든 PPT 생성 완료!")
print(f"저장 위치: {output_dir}")
