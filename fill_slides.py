# -*- coding: utf-8 -*-
"""
1주차.pptx에 Part 2~8 내용 슬라이드를 추가한다.
- 슬라이드 1~11 (Part 1 완성본) 유지
- 슬라이드 12~32 삭제
- Part 2~8: 섹션 구분자 + 내용 슬라이드 추가
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy

SRC = "d:/[1]수업자료/1주차.pptx"
DST = "d:/[1]수업자료/1주차_완성.pptx"

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x3D, 0x3A, 0x35)   # accent4 (Part 라벨 배경)
RED    = RGBColor(0xFF, 0x00, 0x00)   # 강조 (빨강)
BLUE   = RGBColor(0x00, 0x20, 0x60)   # 강조 (짙은 파랑)

# ── 유틸 ────────────────────────────────────────────────────────────────

def delete_slide(prs, idx):
    """슬라이드 삭제 (인덱스 기준)"""
    xml_slides = prs.slides._sldIdLst
    prs.part.drop_rel(xml_slides[idx].rId)
    del xml_slides[idx]

def clone_section_divider(prs, source_idx, part_num, title):
    """슬라이드 복제 후 Part 번호·제목 업데이트 (배경 이미지 포함)"""
    source = prs.slides[source_idx]
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)

    # 이미지 관계 복사 (r:embed 매핑)
    rId_map = {}
    for rId, rel in source.part.rels.items():
        if 'image' in rel.reltype:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    for child in source.shapes._spTree:
        new_child = deepcopy(child)
        for blip in new_child.iter(qn('a:blip')):
            old = blip.get(qn('r:embed'))
            if old and old in rId_map:
                blip.set(qn('r:embed'), rId_map[old])
        sp_tree.append(new_child)

    # 텍스트 업데이트
    for shape in new_slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t.startswith("Part"):
            # "Part X" 텍스트 찾아서 교체
            _set_tf_text(shape.text_frame, f"Part {part_num}", keep_fmt=True)
        elif t in ("제목을 입력하세요", "LLM이란 무엇인가?", "") and \
             shape.left > Inches(2):
            _set_tf_text(shape.text_frame, title, keep_fmt=True)

    return new_slide

def _set_tf_text(tf, new_text, keep_fmt=False):
    """텍스트 프레임의 첫 번째 단락 텍스트를 교체 (서식 유지)"""
    para = tf.paragraphs[0]
    # 기존 런에서 서식 정보 추출
    ref_run = None
    for r in para.runs:
        if r.text.strip():
            ref_run = r
            break
    # 기존 런 모두 삭제
    for r in list(para.runs):
        r._r.getparent().remove(r._r)
    # 새 런 추가
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    from lxml import etree
    rPr_xml = '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    if ref_run is not None:
        try:
            rPr_xml = etree.tostring(ref_run._r.find(qn('a:rPr'))).decode()
        except:
            pass
    import html
    safe_text = html.escape(new_text)
    r_elem = parse_xml(
        f'<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'{rPr_xml}'
        f'<a:t>{safe_text}</a:t></a:r>'
    )
    para._p.append(r_elem)

def add_part_label(slide, part_num):
    """왼쪽 상단 'Part X' 둥근 사각형 라벨 추가"""
    # 둥근 사각형 (채우기: 짙은 갈색)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    rect = slide.shapes.add_shape(
        5,  # msoShapeRoundedRectangle
        Inches(0.40), Inches(0.37), Inches(1.26), Inches(0.52)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK
    rect.line.fill.background()

    # "Part X" 텍스트 (흰색, 가운데 정렬)
    tb = slide.shapes.add_textbox(Inches(0.62), Inches(0.41), Inches(0.82), Inches(0.44))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Part {part_num}"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

def add_title_box(slide, title_text, width_inch=10.5):
    """슬라이드 상단 제목 텍스트 박스 추가"""
    tb = slide.shapes.add_textbox(Inches(1.88), Inches(0.37), Inches(width_inch), Inches(0.55))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True

def add_body(slide, intro, paragraphs):
    """슬라이드 본문 텍스트 추가
    intro: 16pt bold 한 줄 소개문
    paragraphs: list of (text, bold, color) 튜플
    """
    # 소개문 (16pt bold, 맑은 고딕)
    tb_intro = slide.shapes.add_textbox(
        Inches(0.62), Inches(1.10), Inches(11.93), Inches(0.42))
    tf = tb_intro.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = intro
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.name = "맑은 고딕"

    # 본문 (14pt, 여러 단락)
    tb_body = slide.shapes.add_textbox(
        Inches(0.62), Inches(1.65), Inches(11.93), Inches(5.50))
    tf = tb_body.text_frame
    tf.word_wrap = True

    first = True
    for item in paragraphs:
        if isinstance(item, str):
            text, bold, color = item, False, None
        else:
            text, bold, color = item[0], item[1], item[2] if len(item) > 2 else None

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)

        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

def add_content_slide(prs, part_num, title, intro, paragraphs):
    """Part 콘텐츠 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_part_label(slide, part_num)
    add_title_box(slide, title)
    add_body(slide, intro, paragraphs)
    return slide

# ════════════════════════════════════════════════════════════
#  내용 데이터
# ════════════════════════════════════════════════════════════

PARTS = [

    # ── Part 2: Transformer 아키텍처 기초 ──────────────────────
    {
        "part": 2,
        "section_title": "Transformer 아키텍처 기초",
        "slides": [
            {
                "title": "Transformer의 등장",
                "intro": "Transformer는 2017년 Google의 논문 'Attention Is All You Need'에서 제안한 신경망 구조",
                "body": [
                    ("RNN의 한계 극복", True, BLUE),
                    ("① 장기 기억 소실 — 문장이 길어지면 앞부분 정보가 희석됨", False, None),
                    ("② 병렬 처리 불가 — 앞 단계가 끝나야 다음 단계 실행 → 학습 속도 저하", False, None),
                    ("핵심 혁신: Self-Attention 메커니즘", True, BLUE),
                    ("문장 내 모든 토큰이 '서로의 관련성'을 동시에 계산 → 병렬 처리 가능", False, None),
                    ("어떤 단어가 현재 단어와 얼마나 관련 있는지 가중치로 표현", False, None),
                    ("결과", True, BLUE),
                    ("학습 속도 대폭 향상, 더 긴 문맥 처리 가능", False, None),
                    ("→ GPT, BERT, Gemini, Claude 모두 Transformer 아키텍처 기반", True, RED),
                ],
            },
            {
                "title": "인코더와 디코더 구조",
                "intro": "Transformer는 목적에 따라 인코더(이해)와 디코더(생성) 구조로 나뉜다",
                "body": [
                    ("인코더 (Encoder)", True, BLUE),
                    ("입력 문장을 이해하는 역할 — 각 단어의 '의미 벡터(컨텍스트 벡터)' 생성", False, None),
                    ("e.g. BERT : 감성 분석, 분류, 정보 추출 등 '이해' 작업에 특화", True, BLUE),
                    ("디코더 (Decoder)", True, BLUE),
                    ("텍스트를 순차적으로 '생성'하는 역할 — 이전에 출력한 토큰을 참고해 다음 토큰 예측", False, None),
                    ("e.g. GPT 계열, Gemini : 대화·작문·요약 등 '생성' 작업 → 이번 수업의 주역 ★", True, BLUE),
                    ("인코더-디코더 (Encoder-Decoder)", True, BLUE),
                    ("번역·요약처럼 '입력 이해 + 출력 생성'이 모두 필요한 작업", False, None),
                    ("e.g. T5, BART", False, None),
                ],
            },
            {
                "title": "사전학습(Pre-training)과 파인튜닝(Fine-tuning)",
                "intro": "LLM은 사전학습 → 지시학습(SFT) → RLHF 세 단계로 완성된다",
                "body": [
                    ("사전학습 (Pre-training)", True, BLUE),
                    ("방대한 인터넷·책 텍스트로 '다음 토큰 예측' 반복 훈련", False, None),
                    ("언어의 패턴, 지식, 문법, 상식을 자연스럽게 학습", False, None),
                    ("파인튜닝 / 지시학습 (Supervised Fine-Tuning, SFT)", True, BLUE),
                    ("사람이 작성한 '좋은 질문-답변' 예시로 추가 학습", False, None),
                    ("명령어를 따르는 능력, 원하는 형식으로 답하는 능력 배양", False, None),
                    ("RLHF (인간 피드백 강화학습)", True, BLUE),
                    ("사람이 여러 답변 중 더 좋은 것을 선택 → 강화학습으로 반영", False, None),
                    ("→ 사전학습 = 지식 습득 / 파인튜닝 = 행동 교정", True, RED),
                ],
            },
            {
                "title": "LLM의 규모와 창발(Emergence) 현상",
                "intro": "모델 파라미터 수가 일정 규모를 넘으면 예상치 못한 새 능력이 갑자기 나타난다",
                "body": [
                    ("파라미터 = 뇌의 시냅스", True, BLUE),
                    ("수십억~수조 개의 숫자로 언어 패턴을 저장", False, None),
                    ("스케일링 법칙 (Scaling Law)", True, BLUE),
                    ("데이터 · 파라미터 · 컴퓨팅 세 가지가 늘수록 성능이 예측 가능하게 향상", False, None),
                    ("창발 (Emergence)", True, BLUE),
                    ("특정 규모 이상에서 갑자기 나타나는 새 능력 — 산술 계산, 코딩, 언어 번역 등", False, None),
                    ("소규모 모델에는 없던 능력이 대규모 모델에서 갑자기 출현", False, None),
                    ("대표 모델 규모", True, BLUE),
                    ("GPT-3 : 175B / Gemini 1.5 Pro : 수백B (추정) / GPT-4 : ~1T (추정)", False, None),
                ],
            },
        ],
    },

    # ── Part 3: 토큰과 토크나이제이션 ──────────────────────────
    {
        "part": 3,
        "section_title": "토큰과 토크나이제이션",
        "slides": [
            {
                "title": "토큰(Token)이란?",
                "intro": "LLM이 텍스트를 처리하는 최소 단위 — 단어, 서브워드, 또는 문자",
                "body": [
                    ("토크나이저 (Tokenizer)", True, BLUE),
                    ("텍스트를 토큰으로 분해하는 도구 — 모델마다 고유한 토크나이저 보유", False, None),
                    ("BPE (Byte-Pair Encoding)", True, BLUE),
                    ("가장 자주 등장하는 문자 쌍을 반복 병합해 어휘 사전(Vocabulary) 구성", False, None),
                    ("e.g.  'unhappiness' → ['un', 'happi', 'ness']  (3 토큰)", True, BLUE),
                    ("e.g.  'hello world' → ['hello', ' world']  (2 토큰)", True, BLUE),
                    ("한국어 토큰화", True, BLUE),
                    ("교착어 특성 — 조사·어미가 붙는 형태소 단위로 분리됨", False, None),
                    ("같은 내용이라도 영어보다 약 1.5~2배 더 많은 토큰 소비", False, None),
                    ("실습 도구: platform.openai.com/tokenizer 에서 직접 확인 가능", False, None),
                ],
            },
            {
                "title": "컨텍스트 윈도우(Context Window)",
                "intro": "모델이 한 번의 추론에서 처리할 수 있는 최대 토큰 수의 한계",
                "body": [
                    ("컨텍스트 윈도우 = 모델의 '작업 기억 (Working Memory)'", True, BLUE),
                    ("이 범위 안의 모든 토큰이 Self-Attention으로 서로 연산됨", False, None),
                    ("초과 시", True, BLUE),
                    ("오래된 토큰이 자동으로 잘려나가 대화 맥락을 잃게 됨", False, None),
                    ("모델별 컨텍스트 한도", True, BLUE),
                    ("GPT-4o : 128K 토큰", False, None),
                    ("Gemini 1.5 Pro : 1,000K (1M) 토큰 — 업계 최대", False, None),
                    ("Claude 3.5 Sonnet : 200K 토큰", False, None),
                    ("토큰 비용 (Gemini Flash 기준)", True, BLUE),
                    ("입력: $0.075 / 1M tokens   출력: $0.30 / 1M tokens", False, None),
                    ("→ 4주차 '대화 이력 관리'와 직결 — 히스토리 길이 제한이 중요한 이유", True, RED),
                ],
            },
            {
                "title": "토큰과 비용의 관계 — 실전 팁",
                "intro": "API를 사용할 때 토큰 수를 최소화하면 비용과 속도를 동시에 개선할 수 있다",
                "body": [
                    ("Gemini API에서 토큰 수 확인하는 방법", True, BLUE),
                    ("response.usage_metadata.prompt_token_count  (입력 토큰)", False, None),
                    ("response.usage_metadata.candidates_token_count  (출력 토큰)", False, None),
                    ("model.count_tokens('텍스트')  — 요청 전 미리 계산", False, None),
                    ("토큰 절약 전략", True, BLUE),
                    ("① 불필요한 수식어 제거 — 간결한 프롬프트 작성", False, None),
                    ("② 출력 길이 제한 — max_output_tokens 파라미터 설정", False, None),
                    ("③ 히스토리 압축 — 오래된 대화를 요약으로 대체 (4주차 상세 학습)", False, None),
                    ("→ 무료 티어 (Flash 기준): 분당 15 req / 일 1,500 req — 수업용으로 충분", True, BLUE),
                ],
            },
        ],
    },

    # ── Part 4: Gemini API 소개 및 특징 ────────────────────────
    {
        "part": 4,
        "section_title": "Gemini API 소개 및 특징",
        "slides": [
            {
                "title": "Google AI Studio — 플레이그라운드 & API Key 발급",
                "intro": "aistudio.google.com — 구글 계정 하나로 즉시 시작할 수 있는 Gemini API 공식 플랫폼",
                "body": [
                    ("플레이그라운드 (Playground)", True, BLUE),
                    ("코드 없이 브라우저에서 바로 Gemini와 대화 가능", False, None),
                    ("프롬프트 실험, 모델 비교, 파라미터 조정을 시각적으로 수행", False, None),
                    ("API Key 발급 절차", True, BLUE),
                    ("① aistudio.google.com 접속 → 구글 계정 로그인", False, None),
                    ("② 왼쪽 메뉴 'Get API key' 클릭", False, None),
                    ("③ 'Create API key in new project' 클릭", False, None),
                    ("④ 생성된 API Key (AIza로 시작) 즉시 복사 후 저장", False, None),
                    ("⑤ 창을 닫으면 재확인 불가 — 즉시 .env 파일에 붙여넣기", False, None),
                    ("참고: 계정 설정에서 언제든 새 Key 생성 또는 삭제 가능", False, None),
                ],
            },
            {
                "title": "Gemini 모델 라인업 비교",
                "intro": "용도와 비용에 맞는 모델 선택이 중요하다 — 이번 수업은 Flash 기본 사용 ★",
                "body": [
                    ("gemini-1.5-flash  ★ 이번 수업 기본 모델", True, BLUE),
                    ("빠르고 저렴, 멀티모달 지원 / 무료: 분당 15 req · 일 1,500 req", False, None),
                    ("gemini-1.5-pro", True, BLUE),
                    ("고성능, 복잡한 추론과 긴 문서 처리에 적합 / 무료: 분당 2 req · 일 50 req", False, None),
                    ("gemini-1.5-flash-8b", True, BLUE),
                    ("초경량 모델 — 비용 최소화, 단순 분류·요약 작업용", False, None),
                    ("gemini-2.0-flash", True, BLUE),
                    ("최신 버전 — 빠른 속도 + 개선된 코딩·추론 성능", False, None),
                    ("모델 선택 전략", True, BLUE),
                    ("개발·테스트 → Flash  /  복잡한 추론 → Pro  /  비용 절감 → Flash-8B", False, None),
                ],
            },
            {
                "title": "Gemini API 주요 기능",
                "intro": "단순 텍스트 생성을 넘어 멀티모달·스트리밍·Function Calling 등 다양한 기능 제공",
                "body": [
                    ("멀티모달 (Multimodal)", True, BLUE),
                    ("텍스트 + 이미지 + 음성 + 비디오를 함께 입력으로 처리", False, None),
                    ("e.g.  이미지 설명 생성, 차트 분석, 화면 이해", False, None),
                    ("스트리밍 (Streaming)", True, BLUE),
                    ("생성 중인 텍스트를 실시간으로 토큰 단위 전송 → ChatGPT와 동일한 타이핑 효과", False, None),
                    ("Function Calling (7주차)", True, BLUE),
                    ("LLM이 외부 함수·API를 직접 호출하도록 요청하는 기능", False, None),
                    ("JSON 모드 (4주차)", True, BLUE),
                    ("response_mime_type='application/json' → 구조화된 JSON만 반환 강제", False, None),
                    ("임베딩 (5주차)", True, BLUE),
                    ("text-embedding-004 모델로 텍스트를 의미 벡터로 변환 — RAG의 핵심", False, None),
                ],
            },
        ],
    },

    # ── Part 5: 개발 환경 세팅 ─────────────────────────────────
    {
        "part": 5,
        "section_title": "개발 환경 세팅 (Python · venv)",
        "slides": [
            {
                "title": "필수 설치 항목",
                "intro": "수업 실습을 위해 아래 세 가지를 반드시 설치한다",
                "body": [
                    ("Python 3.10 이상", True, BLUE),
                    ("python.org → Downloads → 설치 실행", False, None),
                    ("Windows: 설치 화면에서 'Add Python to PATH' 옵션 체크 필수 !", True, RED),
                    ("Mac: Homebrew 설치 후  brew install python  권장", False, None),
                    ("확인 명령어: python --version  또는  python3 --version", False, None),
                    ("VS Code (Visual Studio Code)", True, BLUE),
                    ("code.visualstudio.com → 설치 후 Python 익스텐션 (ms-python.python) 추가", False, None),
                    ("Git", True, BLUE),
                    ("git-scm.com → 설치 / 확인: git --version", False, None),
                    ("향후 GitHub 연동, 버전 관리, 협업에 필수", False, None),
                ],
            },
            {
                "title": "가상환경(Virtual Environment) — 왜 필요한가?",
                "intro": "프로젝트마다 독립된 파이썬 환경을 만들어 패키지 충돌을 방지한다",
                "body": [
                    ("가상환경이 필요한 이유", True, BLUE),
                    ("프로젝트 A: requests v2.28  /  프로젝트 B: requests v2.31 → 충돌 발생", False, None),
                    ("가상환경은 프로젝트별 독립 패키지 공간을 제공해 충돌을 원천 차단", False, None),
                    ("생성", True, BLUE),
                    ("python -m venv .venv   (현재 폴더에 .venv 폴더 생성)", False, None),
                    ("활성화", True, BLUE),
                    ("Windows:    .venv\\Scripts\\activate", False, None),
                    ("Mac/Linux:  source .venv/bin/activate", False, None),
                    ("확인: 터미널 프롬프트 앞에 (.venv) 표시 여부", False, None),
                    ("비활성화:   deactivate", False, None),
                    ("→ 주의: VS Code에서 인터프리터 선택 시 .venv 내 python 선택 필요", True, RED),
                ],
            },
            {
                "title": "패키지 설치 및 프로젝트 구조",
                "intro": "가상환경 활성화 후 필요한 패키지를 설치하고 프로젝트 폴더를 정리한다",
                "body": [
                    ("패키지 설치 (가상환경 활성화 상태에서)", True, BLUE),
                    ("pip install google-generativeai python-dotenv", False, None),
                    ("설치 확인: pip list  (설치된 패키지 목록 출력)", False, None),
                    ("requirements.txt — 협업을 위한 패키지 목록 관리", True, BLUE),
                    ("pip freeze > requirements.txt       (현재 환경 저장)", False, None),
                    ("pip install -r requirements.txt     (다른 PC에서 재설치)", False, None),
                    ("권장 프로젝트 폴더 구조", True, BLUE),
                    ("my_project/", False, None),
                    ("├── .venv/           ← 가상환경 (GitHub 업로드 금지)", False, None),
                    ("├── .env             ← API Key (GitHub 업로드 금지)", False, None),
                    ("├── .gitignore       ← .venv/ 과 .env 를 반드시 추가", False, None),
                    ("├── main.py          ← 메인 코드", False, None),
                    ("└── requirements.txt ← 패키지 목록", False, None),
                ],
            },
        ],
    },

    # ── Part 6: API Key 발급 및 보안 ───────────────────────────
    {
        "part": 6,
        "section_title": "API Key 발급 및 보안",
        "slides": [
            {
                "title": "API Key 발급 단계별 안내",
                "intro": "Google AI Studio에서 API Key를 발급받는 단계별 절차",
                "body": [
                    ("발급 절차", True, BLUE),
                    ("1단계: aistudio.google.com 접속 → 구글 계정으로 로그인", False, None),
                    ("2단계: 왼쪽 사이드바 'Get API key' 클릭", False, None),
                    ("3단계: 'Create API key in new project' 버튼 클릭", False, None),
                    ("4단계: 생성된 Key (AIza...로 시작) 복사", False, None),
                    ("5단계: 즉시 .env 파일에 붙여넣기 — 창 닫으면 재확인 불가 !", True, RED),
                    ("관리 팁", True, BLUE),
                    ("같은 화면에서 기존 Key 삭제 · 새 Key 생성 가능", False, None),
                    ("프로젝트별로 Key를 분리하면 유출 피해 범위를 최소화할 수 있음", False, None),
                    ("무료 사용 한도 (Flash): 분당 15회 요청 / 하루 1,500회 요청", False, None),
                ],
            },
            {
                "title": "API Key 보안 관리 — 절대 원칙",
                "intro": "API Key는 비밀번호와 같다 — 외부에 노출되면 타인이 무제한 사용·과금 가능",
                "body": [
                    ("안전한 방법 1: .env 파일 사용", True, BLUE),
                    ("프로젝트 루트에 .env 파일 생성 후  GOOGLE_API_KEY=AIza...  형식으로 저장", False, None),
                    ("코드에서: from dotenv import load_dotenv; load_dotenv()", False, None),
                    ("안전한 방법 2: .gitignore 설정", True, BLUE),
                    (".gitignore 파일에 .env 한 줄 추가 → GitHub에 절대 올라가지 않음", False, None),
                    ("절대 하지 말 것", True, RED),
                    ("코드 파일에 API Key 직접 붙여넣기", False, None),
                    ("카카오톡 · SNS · 이메일로 공유", False, None),
                    ("GitHub public repository에 .env 파일 커밋", False, None),
                    ("유출 시 대처", True, BLUE),
                    ("→ 즉시 Google AI Studio에서 해당 Key 삭제 → 새 Key 발급", False, None),
                ],
            },
        ],
    },

    # ── Part 7: 첫 번째 API 호출 실습 ──────────────────────────
    {
        "part": 7,
        "section_title": "첫 번째 API 호출 실습",
        "slides": [
            {
                "title": "Gemini API 호출 코드 구조 이해",
                "intro": "API 호출 코드는 다섯 단계로 구성된다 — 각 단계의 역할을 이해하자",
                "body": [
                    ("1단계: 라이브러리 가져오기 (Import)", True, BLUE),
                    ("import google.generativeai as genai", False, None),
                    ("import os   /   from dotenv import load_dotenv", False, None),
                    ("2단계: API Key 인증 (Configure)", True, BLUE),
                    ("load_dotenv()   →   genai.configure(api_key=os.environ['GOOGLE_API_KEY'])", False, None),
                    ("3단계: 모델 초기화 (Model)", True, BLUE),
                    ("model = genai.GenerativeModel('gemini-1.5-flash')", False, None),
                    ("4단계: 생성 요청 (Generate)", True, BLUE),
                    ("response = model.generate_content('프롬프트 내용')", False, None),
                    ("5단계: 결과 출력 (Output)", True, BLUE),
                    ("print(response.text)    ← 생성된 텍스트", False, None),
                    ("→ SDK가 HTTP 통신, 인증, 재시도, 오류 처리를 모두 대신 처리", True, RED),
                ],
            },
            {
                "title": "실습 — 첫 번째 코드 직접 작성",
                "intro": "아래 코드를 VS Code에서 직접 입력하고 실행해 본다",
                "body": [
                    ("파일: main.py  (가상환경 활성화 상태에서 실행)", True, BLUE),
                    ("import google.generativeai as genai", False, None),
                    ("import os", False, None),
                    ("from dotenv import load_dotenv", False, None),
                    ("", False, None),
                    ("load_dotenv()                    # .env 파일 로드", False, None),
                    ("genai.configure(api_key=os.environ['GOOGLE_API_KEY'])", False, None),
                    ("", False, None),
                    ("model = genai.GenerativeModel('gemini-1.5-flash')", False, None),
                    ("response = model.generate_content('한국의 수도는 어디인가요?')", False, None),
                    ("", False, None),
                    ("print(response.text)", False, None),
                    ("print('입력 토큰:', response.usage_metadata.prompt_token_count)", False, None),
                    ("print('출력 토큰:', response.usage_metadata.candidates_token_count)", False, None),
                    ("실행 방법:  python main.py", True, BLUE),
                ],
            },
            {
                "title": "응답(Response) 객체 완전 이해",
                "intro": "generate_content()가 반환하는 response 객체의 주요 속성을 알아본다",
                "body": [
                    ("response.text", True, BLUE),
                    ("생성된 텍스트 문자열 — 가장 자주 사용하는 속성", False, None),
                    ("response.candidates", True, BLUE),
                    ("생성 후보 목록 (기본값: 1개) — candidates[0].content.parts[0].text 로도 접근 가능", False, None),
                    ("response.usage_metadata", True, BLUE),
                    ("prompt_token_count   — 입력에 사용된 토큰 수", False, None),
                    ("candidates_token_count — 출력에 사용된 토큰 수", False, None),
                    ("response.prompt_feedback", True, BLUE),
                    ("안전 필터에 의한 차단 여부 확인 — block_reason 속성", False, None),
                    ("주요 오류 상황", True, BLUE),
                    ("API Key 오류(401) / 할당량 초과(429) / 안전 필터 차단", False, None),
                    ("→ try-except로 오류를 처리하는 습관 필수", True, RED),
                ],
            },
        ],
    },

    # ── Part 8: Q&A 및 과제 안내 ───────────────────────────────
    {
        "part": 8,
        "section_title": "Q&A 및 과제 안내",
        "slides": [
            {
                "title": "Q&A — 오늘 수업 질문 시간",
                "intro": "오늘 배운 내용 중 이해가 어렵거나 궁금한 점을 자유롭게 질문하세요",
                "body": [
                    ("자주 나오는 질문 예시", True, BLUE),
                    ("Q: API Key를 발급했는데 '401 Invalid API Key' 오류가 나요", False, None),
                    ("   → .env 파일 저장 확인 / load_dotenv() 호출 위치 확인", False, None),
                    ("Q: 응답이 매번 다르게 나와요. 고정할 수 없나요?", False, None),
                    ("   → temperature=0.0 으로 설정 (2주차에 자세히 학습)", False, None),
                    ("Q: 한국어로 질문하면 영어로 답변해요", False, None),
                    ("   → 프롬프트에 '반드시 한국어로 답하세요' 추가 (3주차 내용)", False, None),
                    ("Q: 가상환경이 매번 비활성화됩니다", False, None),
                    ("   → VS Code 터미널 열 때마다 활성화 명령어 실행 필요", False, None),
                    ("   → VS Code 설정에서 Python 인터프리터를 .venv로 지정하면 자동 활성화", False, None),
                    ("그 외 궁금한 점은 언제든지 질문하세요!", True, BLUE),
                ],
            },
            {
                "title": "이번 주 과제 안내",
                "intro": "Gemini API를 직접 호출해 결과를 확인하고 제출한다",
                "body": [
                    ("과제 1: 나의 첫 AI 대화 (필수)", True, BLUE),
                    ("Gemini API를 호출해 아무 질문이나 하고, 응답 출력 화면을 캡처해 제출", False, None),
                    ("과제 2: 토큰 수 확인 (필수)", True, BLUE),
                    ("response.usage_metadata로 입력 · 출력 토큰 수를 출력하고 스크린샷 제출", False, None),
                    ("과제 3: 자유 탐구 (선택)", True, BLUE),
                    ("temperature 값을 0.0, 1.0, 2.0으로 각각 바꿔보고 응답 차이 비교 제출", False, None),
                    ("제출 방법 및 기한", True, BLUE),
                    ("제출처: 학교 LMS (eClass) 과제 게시판", False, None),
                    ("기한: 다음 수업 전날 자정까지", False, None),
                    ("다음 주 예고", True, BLUE),
                    ("Python SDK 심화 + GenerationConfig + CLI 챗봇 구현 실습", False, None),
                ],
            },
        ],
    },
]

# ════════════════════════════════════════════════════════════
#  메인 실행
# ════════════════════════════════════════════════════════════

def main():
    prs = Presentation(SRC)
    total = len(prs.slides)
    print(f"원본 슬라이드 수: {total}")

    # 1. 슬라이드 12~32 (인덱스 11~31) 삭제 — 뒤에서부터
    for i in range(total - 1, 10, -1):
        delete_slide(prs, i)
    print(f"삭제 후 슬라이드 수: {len(prs.slides)}  (기대: 11)")

    # 2. 섹션 구분자(슬라이드 2, 인덱스 2) 클론 함수 준비
    SECTION_DIVIDER_IDX = 2   # "Part 1 | LLM이란 무엇인가?" 슬라이드

    # 3. Part 2~8 추가
    for part_data in PARTS:
        pnum  = part_data["part"]
        stitle = part_data["section_title"]
        print(f"\n  Part {pnum}: {stitle}")

        # 섹션 구분자 복제 및 텍스트 업데이트
        clone_section_divider(prs, SECTION_DIVIDER_IDX, pnum, stitle)
        print(f"    섹션 구분자 추가")

        # 콘텐츠 슬라이드
        for slide_data in part_data["slides"]:
            add_content_slide(
                prs,
                part_num=pnum,
                title=slide_data["title"],
                intro=slide_data["intro"],
                paragraphs=slide_data["body"],
            )
            print(f"    콘텐츠 슬라이드 추가: {slide_data['title'][:30]}")

    print(f"\n최종 슬라이드 수: {len(prs.slides)}")
    prs.save(DST)
    print(f"저장 완료: {DST}")

if __name__ == "__main__":
    main()
